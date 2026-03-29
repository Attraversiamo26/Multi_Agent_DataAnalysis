"""
Data Agent Streamlit Frontend
A beautiful chat interface for the Data Agent API
"""

import streamlit as st
import requests
import json
import logging
from datetime import datetime
from typing import Generator, Dict, Any
import time
import re
import ast
import os
import pandas as pd

logger = logging.getLogger(__name__)

# Page configuration
st.set_page_config(
    page_title="牛马小熊猫智能体",
    page_icon="📮",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI - inspired by reference design
st.markdown("""
<style>
    /* Main container */
    .main-container {
        display: flex;
        height: 100vh;
        background-color: #f5f5f5;
    }
    
    /* Sidebar navigation */
    .sidebar {
        width: 200px;
        background-color: white;
        border-right: 1px solid #e0e0e0;
        padding: 20px;
        box-shadow: 2px 0 5px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .sidebar-header {
        display: flex;
        align-items: center;
        margin-bottom: 30px;
        padding-bottom: 15px;
        border-bottom: 1px solid #f0f0f0;
    }
    
    .sidebar-logo {
        width: 100px;
        height: 100px;
        margin-right: 2px;
        object-fit: contain;
    }
    
    .sidebar-title {
        font-size: 20px;
        font-weight: 600;
        color: black;
        line-height: 1.2;
    }
    
    .nav-item {
        display: flex;
        align-items: center;
        padding: 12px 15px;
        margin: 5px 0;
        border-radius: 8px;
        cursor: pointer;
        transition: all 0.3s ease;
        color: #666;
        text-decoration: none;
    }
    
    .nav-item:hover {
        background-color: #f0f8ff;
        color: #1a73e8;
        transform: translateX(5px);
    }
    
    .nav-item.active {
        background-color: #e3f2fd;
        color: #1a73e8;
        font-weight: 500;
        box-shadow: 0 2px 8px rgba(26, 115, 232, 0.1);
    }
    
    .nav-icon {
        margin-right: 10px;
        font-size: 18px;
    }
    
    /* Main content area */
    .main-content {
        flex: 1;
        display: flex;
        flex-direction: column;
        overflow: hidden;
    }
    
    /* Top bar */
    .top-bar {
        background-color: white;
        padding: 15px 30px;
        border-bottom: 1px solid #e0e0e0;
        display: flex;
        justify-content: space-between;
        align-items: center;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    
    .search-box {
        display: flex;
        align-items: center;
        background-color: #f5f5f5;
        border-radius: 20px;
        padding: 8px 15px;
        width: 300px;
        transition: all 0.3s ease;
    }
    
    .search-box:hover {
        background-color: #e8e8e8;
    }
    
    .search-box input {
        border: none;
        background: transparent;
        outline: none;
        flex: 1;
        margin-left: 10px;
    }
    
    .user-profile {
        display: flex;
        align-items: center;
    }
    
    .user-avatar {
        width: 36px;
        height: 36px;
        border-radius: 50%;
        background-color: #e3f2fd;
        display: flex;
        align-items: center;
        justify-content: center;
        margin-left: 15px;
        transition: all 0.3s ease;
    }
    
    .user-avatar:hover {
        transform: scale(1.1);
        box-shadow: 0 2px 8px rgba(26, 115, 232, 0.2);
    }
    
    /* Welcome section */
    .welcome-section {
        text-align: center;
        padding: 40px 20px;
        background-color: #fafafa;
        flex: 1;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
    }
    
    .welcome-title {
        font-size: 28px;
        font-weight: 600;
        color: #333;
        margin-bottom: 20px;
        animation: fadeInUp 0.8s ease;
    }
    
    .welcome-subtitle {
        font-size: 16px;
        color: #666;
        margin-bottom: 40px;
        animation: fadeInUp 1s ease;
    }
    
    .logo-container {
        margin: 20px 0;
        animation: fadeInUp 1.2s ease;
    }
    
    .main-logo {
        width: 200px;
        height: 200px;
        object-fit: contain;
    }
    
    /* Background header */
    .background-header {
        background-image: url('frontend_assets/抬头图片.png');
        background-size: cover;
        background-position: center;
        border-radius: 12px;
        padding: 40px;
        width: 100%;
        max-width: 800px;
        margin: 0 auto 20px auto;
        text-align: center;
        color: white;
        box-shadow: 0 4px 12px rgba(0,0,0,0.2);
        animation: fadeIn 1s ease;
    }
    
    .header-title {
        font-size: 24px;
        font-weight: 600;
        margin: 0;
    }
    
    /* Upload area */
    .upload-area {
        background-color: white;
        border-radius: 12px;
        padding: 30px;
        width: 100%;
        max-width: 600px;
        margin: 0 auto;
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        color: #000000;
        transition: all 0.3s ease;
        animation: fadeInUp 0.8s ease;
    }
    
    .upload-area:hover {
        box-shadow: 0 6px 20px rgba(0,0,0,0.15);
        transform: translateY(-2px);
    }
    
    .upload-header {
        display: flex;
        align-items: center;
        margin-bottom: 20px;
    }
    
    .upload-icon {
        width: 40px;
        height: 40px;
        border-radius: 50%;
        background-color: #4CAF50;
        display: flex;
        align-items: center;
        justify-content: center;
        margin-right: 15px;
        animation: pulse 2s infinite;
    }
    
    .upload-title {
        font-size: 18px;
        font-weight: 600;
        color: #000000;
        margin: 0;
    }
    
    .upload-description {
        background-color: #f5f5f5;
        border-radius: 8px;
        padding: 20px;
        min-height: 100px;
        margin-bottom: 20px;
        text-align: center;
    }
    
    .upload-description p {
        color: #000000;
        margin: 0;
    }
    
    .upload-content {
        display: flex;
        flex-direction: column;
        gap: 15px;
    }
    
    .drag-drop-area {
        border: 2px dashed #4CAF50;
        border-radius: 8px;
        padding: 20px;
        text-align: center;
        transition: all 0.3s ease;
    }
    
    .drag-drop-area:hover {
        background-color: rgba(76, 175, 80, 0.05);
        border-color: #45a049;
    }
    
    .drag-drop-text {
        color: #000000;
        margin: 0;
    }
    
    .input-area {
        display: flex;
        gap: 10px;
    }
    
    .input-field {
        flex: 1;
        padding: 15px;
        border: 1px solid #e0e0e0;
        border-radius: 20px;
        outline: none;
        color: #000000;
        transition: all 0.3s ease;
    }
    
    .input-field:focus {
        border-color: #4CAF50;
        box-shadow: 0 0 0 2px rgba(76, 175, 80, 0.2);
    }
    
    .input-field::placeholder {
        color: #666666;
    }
    
    .send-button {
        background-color: #4CAF50;
        color: white;
        border: none;
        border-radius: 50%;
        width: 40px;
        height: 40px;
        cursor: pointer;
        display: flex;
        align-items: center;
        justify-content: center;
        transition: all 0.3s ease;
    }
    
    .send-button:hover {
        background-color: #45a049;
        transform: scale(1.1);
        box-shadow: 0 2px 8px rgba(76, 175, 80, 0.3);
    }
    
    /* Chat area */
    .chat-area {
        flex: 1;
        display: flex;
        flex-direction: column;
        padding: 20px;
        overflow: hidden;
    }
    
    .chat-messages {
        flex: 1;
        overflow-y: auto;
        margin-bottom: 20px;
        padding-right: 10px;
    }
    
    /* Custom scrollbar */
    .chat-messages::-webkit-scrollbar {
        width: 6px;
    }
    
    .chat-messages::-webkit-scrollbar-track {
        background: #f1f1f1;
        border-radius: 3px;
    }
    
    .chat-messages::-webkit-scrollbar-thumb {
        background: #c1c1c1;
        border-radius: 3px;
    }
    
    .chat-messages::-webkit-scrollbar-thumb:hover {
        background: #a8a8a8;
    }
    
    .chat-input-area {
        display: flex;
        gap: 10px;
        padding: 10px;
        background-color: white;
        border-radius: 10px;
        box-shadow: 0 -2px 10px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .chat-input {
        flex: 1;
        border: 1px solid #e0e0e0;
        border-radius: 20px;
        padding: 12px 20px;
        outline: none;
        transition: all 0.3s ease;
    }
    
    .chat-input:focus {
        border-color: #1a73e8;
        box-shadow: 0 0 0 2px rgba(26, 115, 232, 0.2);
    }
    
    /* Message styles */
    .message {
        margin-bottom: 15px;
        max-width: 70%;
        animation: fadeIn 0.5s ease;
    }
    
    .user-message {
        margin-left: auto;
        background-color: #e3f2fd;
        border-radius: 18px 18px 4px 18px;
        padding: 12px 16px;
        box-shadow: 0 2px 4px rgba(26, 115, 232, 0.1);
    }
    
    .assistant-message {
        margin-right: auto;
        background-color: white;
        border-radius: 18px 18px 18px 4px;
        padding: 12px 16px;
        border: 1px solid #e0e0e0;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    
    /* Data preview */
    .data-preview {
        background-color: white;
        border-radius: 10px;
        padding: 20px;
        margin-top: 20px;
        box-shadow: 0 2px 10px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .data-preview:hover {
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }
    
    /* Button styles */
    .stButton > button {
        transition: all 0.3s ease;
        border-radius: 8px !important;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    /* Card styles */
    .stCard {
        background-color: white;
        border-radius: 12px;
        padding: 20px;
        box-shadow: 0 2px 10px rgba(0,0,0,0.05);
        margin-bottom: 20px;
        transition: all 0.3s ease;
    }
    
    .stCard:hover {
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        transform: translateY(-2px);
    }
    
    /* Animations */
    @keyframes fadeIn {
        from { opacity: 0; }
        to { opacity: 1; }
    }
    
    @keyframes fadeInUp {
        from {
            opacity: 0;
            transform: translateY(20px);
        }
        to {
            opacity: 1;
            transform: translateY(0);
        }
    }
    
    @keyframes pulse {
        0% {
            box-shadow: 0 0 0 0 rgba(76, 175, 80, 0.4);
        }
        70% {
            box-shadow: 0 0 0 10px rgba(76, 175, 80, 0);
        }
        100% {
            box-shadow: 0 0 0 0 rgba(76, 175, 80, 0);
        }
    }
    
    /* Responsive adjustments */
    @media (max-width: 768px) {
        .sidebar {
            width: 160px;
            padding: 15px;
        }
        
        .sidebar-title {
            font-size: 16px;
        }
        
        .upload-area {
            width: 95%;
            padding: 20px;
        }
        
        .background-header {
            width: 95%;
            padding: 30px;
        }
        
        .welcome-title {
            font-size: 24px;
        }
        
        .chat-messages {
            max-width: 90%;
        }
        
        .search-box {
            width: 200px;
        }
    }
    
    @media (max-width: 480px) {
        .sidebar {
            width: 140px;
        }
        
        .nav-item {
            padding: 10px;
        }
        
        .nav-icon {
            font-size: 16px;
        }
        
        .top-bar {
            padding: 10px 15px;
        }
        
        .chat-area {
            padding: 10px;
        }
    }
    
    /* 自定义输入框样式 - 通过key定位 */
    [data-testid="stTextInput"][key="sidebar_search"] {
        background-color: #f5f5f5;
        border-radius: 20px;
        padding: 8px 15px;
        width: 100%;
        border: none;
        transition: all 0.3s ease;
    }
    
    [data-testid="stTextInput"][key="sidebar_search"]:hover {
        background-color: #e8e8e8;
    }
    
    [data-testid="stTextInput"][key="sidebar_search"] input {
        border: none;
        background: transparent;
        outline: none;
        flex: 1;
        margin-left: 10px;
    }
    
    /* DataFrame styles */
    .dataframe {
        border-radius: 8px;
        overflow: hidden;
    }
    
    /* Expander styles */
    .streamlit-expander {
        border-radius: 8px;
        margin-bottom: 10px;
    }
    
    .streamlit-expanderHeader {
        border-radius: 8px;
        transition: all 0.3s ease;
    }
    
    .streamlit-expanderHeader:hover {
        background-color: #f5f5f5;
    }
</style>
""", unsafe_allow_html=True)


class DataAgentClient:
    """Client for Data Agent API"""

    def __init__(self, base_url: str):
        self.base_url = base_url.rstrip('/')
        self.chat_endpoint = f"{self.base_url}/api/chat/stream"
        self.similar_questions_endpoint = f"{self.base_url}/api/similar-questions"

    def check_server_status(self) -> bool:
        """Check if server is running"""
        try:
            response = requests.get(self.base_url, timeout=2)
            return True
        except requests.exceptions.RequestException:
            return False

    def chat_stream(self, user_input: str, session_id: str) -> Generator[Dict[str, Any], None, None]:
        messages = [{
            "role": "user",
            "content": user_input
        }]

        payload = {
            "messages": messages,
            "session_id": session_id
        }

        # 新增：去重缓存，避免重复事件
        processed_events = set()
        
        try:
            response = requests.post(
                self.chat_endpoint,
                json=payload,
                headers={"Content-Type": "application/json"},
                stream=True,
                timeout=180
            )

            if response.status_code != 200:
                yield {"type": "error", "content": f"❌ Error: HTTP {response.status_code}", "id": ""}
                return

            last_message_id = None
            current_event_type = None
            for line in response.iter_lines(decode_unicode=True):
                # 修复：严格过滤空行/无效行
                if not line or line.strip() == "" or line.strip() == ":":
                    continue

                if line.startswith("event:"):
                    current_event_type = line.split(":", 1)[1].strip()
                    continue

                if line.startswith("data:"):
                    data_str = line.split(":", 1)[1].strip()
                    # 修复：跳过空数据
                    if not data_str:
                        continue
                    # 修复：去重，避免重复日志
                    if data_str in processed_events:
                        continue
                    processed_events.add(data_str)
                    
                    try:
                        data = json.loads(data_str)
                        content = data.get("content", "")
                        message_id = data.get("id", "")

                        if isinstance(content, str) and content:
                            if last_message_id and message_id != last_message_id:
                                yield {"type": "separator", "id": message_id, "content": ""}
                            yield {
                                "type": current_event_type or "message",
                                "id": message_id or "",
                                "content": content
                            }
                            if message_id:
                                last_message_id = message_id
                    except json.JSONDecodeError:
                        continue
        except requests.exceptions.Timeout:
            yield {"type": "error", "content": "❌ Request timeout. Please try again.", "id": ""}
        except requests.exceptions.ConnectionError:
            yield {"type": "error", "content": "❌ Cannot connect to server. Please make sure the server is running.", "id": ""}
        except Exception as e:
            yield {"type": "error", "content": f"❌ Error: {str(e)}", "id": ""}

    def get_similar_questions(self, user_input: str, context: list) -> list:
        """
        Get similar questions based on user input and context

        Args:
            user_input: User's current question
            context: List of previous messages for context

        Returns:
            List of similar questions
        """
        try:
            payload = {
                "user_question": user_input,
                "context": context
            }

            response = requests.post(
                self.similar_questions_endpoint,
                json=payload,
                headers={"Content-Type": "application/json"},
                timeout=30
            )

            if response.status_code == 200:
                data = response.json()
                return data.get("similar_questions", [])
            else:
                return []
        except Exception as e:
            print(f"Error getting similar questions: {str(e)}")
            return []

import re

def clean_raw_content(content: str) -> str:
    """清理原始内容中的路径和标记"""
    content = re.sub(r'zh-CN[^\s]*', '', content)
    content = re.sub(r'thread_\d+_\d+[^\s]*', '', content)
    content = re.sub(r'/Users/[^\s]+', '', content)
    content = re.sub(r'\s{3,}', '  ', content)
    return content.strip()

def detect_and_format_code(content: str) -> tuple[str, list[tuple[str, str]]]:
    """
    检测并格式化代码块，包括从run_python_code返回结果中提取代码
    
    返回:
        (处理后的内容, [(语言, 代码), ...])
    """
    code_blocks = []
    
    def replace_and_save(match, lang):
        code_blocks.append((lang, match.group(1)))
        return f'__CODE_BLOCK_{len(code_blocks)-1}__'
    
    # Python代码 (优先)
    python_pattern = r'```python\s*([\s\S]*?)\s*```'
    content = re.sub(python_pattern, lambda m: replace_and_save(m, 'python'), content)
    
    # JSON代码
    json_pattern = r'```json\s*([\s\S]*?)\s*```'
    content = re.sub(json_pattern, lambda m: replace_and_save(m, 'json'), content)
    
    # 通用代码块 - 智能识别
    generic_pattern = r'```\s*([\s\S]*?)\s*```'
    for match in list(re.finditer(generic_pattern, content)):
        if '__CODE_BLOCK_' not in match.group(0):
            code_content = match.group(1)
            if 'import ' in code_content or 'def ' in code_content or 'pandas' in code_content:
                lang = 'python'
            elif code_content.strip().startswith(('{', '[')):
                lang = 'json'
            else:
                lang = 'text'
            code_blocks.append((lang, code_content))
            content = content.replace(match.group(0), f'__CODE_BLOCK_{len(code_blocks)-1}__')
    
    return content, code_blocks

def parse_and_classify_json(content: str) -> dict:
    """
    解析并分类 JSON 数据，按照新的前端展示标准
    
    前端展示模块 | 必填字段（代码解析 + 前端渲染用）
    1. 意图识别 | intent_type, confidence, reasoning
    2. 执行过程 | execution_records(action, agent, execution_status, output)
    3. 结果数据 | result_data (total_count, 统计值，表格数据)
    4. Python 脚本 | python_code
    5. 路由日志 | routing_log
    """
    result = {
        'intent_recognition': None,
        'plan': None,
        'execution_records': [],
        'result_data': [],
        'python_code': [],
        'routing_log': None
    }

    # 1. 提取路由日志
    routing_match = re.search(r'(src\.graph\.enhanced_builder[^\\n]*Routing based on intent[^\\n]*)', content)
    if routing_match:
        result['routing_log'] = routing_match.group(1)

    # 2. 提取所有 JSON 代码块
    json_blocks = []
    for match in re.finditer(r'```json\s*([\s\S]*?)\s*```', content):
        try:
            json_blocks.append(json.loads(match.group(1)))
        except:
            pass

    # 放宽裸露 JSON 提取规则
    json_patterns = [r'(\{[\s\S]*\})', r'(\[[\s\S]*\])']
    for pattern in json_patterns:
        matches = re.findall(pattern, content)
        for match in matches:
            try:
                json_data = json.loads(match)
                json_blocks.append(json_data)
            except:
                pass

    # 3. 按照新标准分类数据
    for json_data in json_blocks:
        if isinstance(json_data, dict):
            # 意图识别
            if 'intent_type' in json_data:
                result['intent_recognition'] = json_data
            # 执行计划
            elif 'steps' in json_data:
                result['plan'] = json_data
            # 执行记录 - 核心修复：确保所有步骤都记录
            elif 'execution_status' in json_data or 'action' in json_data or 'step_name' in json_data:
                result['execution_records'].append(json_data)
            # 结果数据 - 其他有效数据
            else:
                result['result_data'].append(json_data)
        elif isinstance(json_data, list):
            result['result_data'].append(json_data)
    
    # 4. 提取 Python 代码块
    for match in re.finditer(r'```python\s*([\s\S]*?)\s*```', content):
        result['python_code'].append(match.group(1))

    return result

def save_code_to_file(code: str, filename: str, directory: str = "extracted_scripts") -> str:
    """保存代码到文件"""
    os.makedirs(directory, exist_ok=True)
    filepath = os.path.join(directory, filename)
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(code)
    return filepath

def extract_tabular_data_from_results(result_data: list) -> list[tuple[str, pd.DataFrame]]:
    """
    专门从结果数据中提取最后一个步骤的表格化数据
    修复：增加空值、数据类型校验，避免计算报错
    """
    tables = []
    
    for json_data in reversed(result_data):
        try:
            if isinstance(json_data, dict):
                for key in ['result', 'results', 'data', 'records', 'items', 'observation']:
                    if key in json_data:
                        val = json_data[key]
                        # 核心修复：校验数据格式，非空+字典列表
                        if isinstance(val, list) and len(val) > 0 and isinstance(val[0], dict):
                            # 修复：过滤非数值列，避免计算报错
                            df = pd.DataFrame(val)
                            # 清理空列/无效列
                            df = df.dropna(how='all', axis=1)
                            df = df.loc[:, ~df.columns.str.contains('^Unnamed', regex=True)]
                            tables.append((key, df))
                            return tables
            
            elif isinstance(json_data, list) and len(json_data) > 0 and isinstance(json_data[0], dict):
                df = pd.DataFrame(json_data)
                df = df.dropna(how='all', axis=1)
                tables.append(("最终结果", df))
                return tables
        except Exception as e:
            # 修复：捕获计算异常，返回错误提示而非崩溃
            tables.append(("计算错误", pd.DataFrame([{"错误信息": str(e)}])))
            return tables
    
    # 兜底：无数据时返回空提示
    if not tables:
        tables.append(("提示", pd.DataFrame([{"信息": "暂无计算结果"}])))
    return tables

def process_and_display_assistant_message(content: str):
    content = clean_raw_content(content)
    content_with_placeholders, code_blocks = detect_and_format_code(content)
    json_classification = parse_and_classify_json(content)

    # 核心修复：步骤去重，避免重复展示
    unique_records = []
    step_names = set()
    for record in json_classification['execution_records']:
        step_name = record.get('step_name', '')
        if step_name not in step_names and step_name:
            step_names.add(step_name)
            unique_records.append(record)
    json_classification['execution_records'] = unique_records

    python_codes = []
    saved_code_files = []
    
    # 修复：代码块提取增强
    for lang, code in code_blocks:
        if lang == 'python' and code.strip():
            python_codes.append(code)
            filename = f"extracted_code_{len(saved_code_files)+1}.py"
            try:
                saved_path = save_code_to_file(code, filename)
                saved_code_files.append((filename, saved_path))
            except Exception as e:
                logger.warning(f"保存代码文件失败: {e}")

    tables = extract_tabular_data_from_results(json_classification['result_data'])

    with st.expander("💭 执行过程（点击展开）", expanded=False):
        # 意图识别
        if json_classification['intent_recognition'] or json_classification['routing_log']:
            st.markdown("### 意图识别")
            if json_classification['intent_recognition']:
                st.json(json_classification['intent_recognition'], expanded=True)
            if json_classification['routing_log']:
                st.markdown(f"**路由日志：** {json_classification['routing_log']}")
        
        # 执行计划
        if json_classification['plan']:
            st.markdown("### 执行计划")
            st.json(json_classification['plan'], expanded=True)
        
        # 执行过程（已去重）
        if json_classification['execution_records']:
            st.markdown("### 执行过程")
            for idx, record in enumerate(json_classification['execution_records'], 1):
                step_name = record.get('step_name', f'Step {idx}')
                st.markdown(f"**步骤 {idx}: {step_name}**")
                st.json(record, expanded=True)

    # 结果展示
    has_content = tables or python_codes
    if has_content:
        st.markdown("## 📊 结果返回")
        if tables:
            st.markdown("### 表格化数据")
            for table_name, df in tables:
                st.markdown(f"**{table_name}**")
                st.dataframe(df, use_container_width=True)
        # 修复：强制展示Python脚本
        if python_codes:
            st.markdown("### Python 代码脚本")
            for i, code in enumerate(python_codes, 1):
                with st.expander(f"代码 {i}", expanded=True):
                    st.code(code, language='python', line_numbers=True)
                    if i <= len(saved_code_files):
                        filename, filepath = saved_code_files[i-1]
                        st.success(f"✅ 已保存至: {filepath}")

def initialize_session_state():
    """Initialize Streamlit session state"""
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "session_id" not in st.session_state:
        st.session_state.session_id = f"thread_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    if "server_url" not in st.session_state:
        st.session_state.server_url = "http://localhost:10000"
    if "client" not in st.session_state:
        st.session_state.client = DataAgentClient(st.session_state.server_url)
    if "current_page" not in st.session_state:
        st.session_state.current_page = "home"
    if "uploaded_file" not in st.session_state:
        st.session_state.uploaded_file = None

def render_sidebar():
    """Render sidebar with navigation"""
    with st.sidebar:
        col1, col2 = st.columns([1, 3], vertical_alignment="center")
        with col1:
            # 修复图片路径问题：如果文件不存在则显示文本兜底
            try:
                st.image("frontend_assets/公司图标.png", width=80)  # 修正图片宽度（原800过宽）
            except:
                st.markdown('<div style="width:80px;height:80px;background:#e3f2fd;border-radius:8px;display:flex;align-items:center;justify-content:center;">📮</div>', unsafe_allow_html=True)
        with col2:
            st.markdown("""
            <div class='sidebar-title' style='display: flex; align-items: center; height: 100%; line-height: 1.5;'>
                天才小熊猫
            </div>
            """, unsafe_allow_html=True)

        # Navigation items
        nav_items = [
            {"icon": "🏠", "label": "首页", "page": "home"},
            {"icon": "📊", "label": "数据分析", "page": "analysis"},
            {"icon": "🐍", "label": "Python沙盒", "page": "python_sandbox"},
            {"icon": "📄", "label": "我的文档", "page": "documents"},
            {"icon": "🧠", "label": "知识库", "page": "knowledge"},
            {"icon": "📚", "label": "历史会话", "page": "history"}
        ]

        for item in nav_items:
            btn_key = f"nav_btn_{item['page']}"
            if st.button(f"{item['icon']} {item['label']}", width='stretch', key=btn_key):
                st.session_state.current_page = item["page"]
                st.rerun()

        # Session management
        st.markdown("""
        <div style="margin-top: 40px; padding-top: 20px; border-top: 1px solid #f0f0f0;">
            <div style="font-size: 18px; color: #666; margin-bottom: 10px;">分析会话管理</div>
        </div>
        """, unsafe_allow_html=True)

        if st.button("🔄 新会话", width='stretch', key="new_session_btn"):
            st.session_state.messages = []
            st.session_state.session_id = f"thread_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            st.rerun()

        if st.button("🗑️ 清除聊天", width='stretch', key="clear_chat_btn"):
            st.session_state.messages = []
            st.rerun()

        # 搜索框 - 移除class_，改用key定位样式
        st.markdown("### 🔍 快速搜索")
        search_text = st.text_input("搜索内容...", key="sidebar_search")

def render_home_page():
    """Render home page with welcome screen"""
    # Main title with background
    try:
        col_text, col_img = st.columns([2, 1], gap="large")

        # 左侧：文字列
        with col_text:
            st.markdown("""
            <h1 style='font-size: 38px; font-weight: 700; margin: 0; color: #333; line-height: 1.5;'>
                欢迎使用牛马小熊猫智能体<br>——您的上班摸鱼小助手
            </h1>
            <p style='font-size: 16px; color: #666; margin-top: 20px; line-height: 1.6;'>
                支持CSV/Excel数据上传、智能分析、可视化展示,助力工作高效处理与决策
            </p>
            """, unsafe_allow_html=True)

        # 右侧：图片列
        with col_img:
            if os.path.exists('frontend_assets/抬头图片.png'):
                st.image(
                    "frontend_assets/抬头图片.png",
                    width=350,  # 可根据图片实际尺寸调整
                    use_container_width=False
                )
            else:
                # 图片不存在时的兜底展示
                st.markdown("""
                <div style='width:350px;height:200px;background:#4CAF50;border-radius:12px;display:flex;align-items:center;justify-content:center;'>
                    <h2 style='color:white;'>牛马小熊猫智能体</h2>
                </div>
                """, unsafe_allow_html=True)

    except Exception as e:
        # 异常兜底展示
        st.error(f"加载头部内容出错：{str(e)}")
        st.markdown("""
        <div style='background-color: #4CAF50; border-radius: 12px; padding: 60px; width: 100%;
                   max-width: 800px; margin: 0 auto 20px auto; text-align: center;
                   box-shadow: 0 4px 12px rgba(0,0,0,0.2);'>
            <h1 style='font-size: 36px; font-weight: 700; margin: 0; color: white;'>
                欢迎使用牛马小熊猫智能体
            </h1>
        </div>
        """, unsafe_allow_html=True)

    # Custom CSS for better text wrapping, chat display, and collapsible thinking
    st.markdown("""
    <style>
    /* Ensure all text wraps properly */
    .stMarkdown, .chat-message, div[data-testid="stMarkdownContainer"] {
        word-wrap: break-word !important;
        white-space: normal !important;
        overflow-wrap: break-word !important;
        hyphens: auto;
    }
    
    /* Chat message containers */
    .stChatMessage {
        word-wrap: break-word !important;
        white-space: normal !important;
    }
    
    /* Markdown content */
    .stMarkdown p {
        word-wrap: break-word !important;
        white-space: normal !important;
        line-height: 1.6;
    }
    
    /* Code blocks */
    pre {
        white-space: pre-wrap !important;
        word-wrap: break-word !important;
        overflow-wrap: break-word !important;
    }
    
    /* Dataframes */
    .stDataFrame {
        max-width: 100%;
        overflow-x: auto;
    }
    
    /* Collapsible thinking box */
    .thinking-box {
        background-color: #f0f2f6;
        border-left: 4px solid #1f77b4;
        border-radius: 8px;
        padding: 12px 16px;
        margin: 10px 0;
        font-size: 0.9em;
    }
    
    .thinking-box summary {
        cursor: pointer;
        font-weight: 600;
        color: #1f77b4;
        outline: none;
    }
    
    .thinking-box summary:hover {
        color: #145a8a;
    }
    
    .thinking-content {
        margin-top: 10px;
        padding-left: 10px;
        color: #555;
    }
    
    /* Response content */
    .response-content {
        background-color: #ffffff;
        border-radius: 8px;
        padding: 12px 16px;
        margin: 10px 0;
    }
    </style>
    """, unsafe_allow_html=True)
    st.markdown("### 📁 数据上传")
    uploaded_file = st.file_uploader(
        "上传您的数据文件",
        type=["csv", "xlsx", "xls"],
        help="支持的格式：CSV、Excel",
        key="home_file_uploader"
    )

    if uploaded_file is not None:
        # Create csv_files directory if it doesn't exist
        os.makedirs("csv_files", exist_ok=True)

        # Save the uploaded file first
        file_path = os.path.join("csv_files", uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # Then update conf.yaml to add the uploaded file to data_sources
        import subprocess
        result = subprocess.run(["python3", "update_csv_config.py"], capture_output=True, text=True)
        if "Updated" in result.stdout:
            st.success(f"✅ 文件已保存：{uploaded_file.name}，并已添加到配置中")
        else:
            st.success(f"✅ 文件已保存：{uploaded_file.name}")

        st.session_state.uploaded_file = uploaded_file

        # Show file preview
        try:
            if uploaded_file.name.endswith(".csv"):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            df = df.loc[:, ~df.columns.str.contains('^Unnamed', regex=True)]
            st.subheader("📊 文件预览")
            st.dataframe(df.head(), width='stretch')
            st.text(f"数据形状：{df.shape[0]} 行 × {df.shape[1]} 列")
            
            # 以表格形式展示所有列名
            st.subheader("📋 列名信息")
            columns_df = pd.DataFrame({
                '列名': df.columns.tolist(),
                '数据类型': [str(dtype) for dtype in df.dtypes.tolist()]
            })
            st.dataframe(columns_df, width='stretch')
        except Exception as e:
            st.error(f"❌ 读取文件时出错：{str(e)}")

    # Chat section
    st.markdown("### 💬 对话分析")

    # Display chat messages with enhanced formatting
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            content = message["content"]
            
            if message["role"] == "assistant":
                process_and_display_assistant_message(content)
            else:
                st.markdown(content)

    # Chat input
    user_input = st.chat_input("请输入您的分析需求...", key="home_chat_input")

    # Process user input
    if user_input:
        # Check server connection first
        if not st.session_state.client.check_server_status():
            st.error("❌ 无法连接到服务器。请检查服务器URL并确保服务器正在运行。")
            return

        # Add user message to chat
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)

        # Get assistant response
        with st.chat_message("assistant"):
            full_response = ""
            response_placeholder = st.empty()
            
            # 显示正在处理的提示，而不是显示原始内容
            response_placeholder.markdown("⏳ 正在处理您的请求...")
            
            for event in st.session_state.client.chat_stream(
                user_input,
                st.session_state.session_id
            ):
                content_chunk = event.get("content", "")
                if content_chunk:
                    full_response += content_chunk

            # Finalize response - 清空占位符，然后调用处理函数展示
            response_placeholder.empty()
            process_and_display_assistant_message(full_response)
            st.session_state.messages.append({"role": "assistant", "content": full_response})

        # Get similar questions after receiving response
        similar_questions = st.session_state.client.get_similar_questions(
            user_input, 
            st.session_state.messages[-5:]
        )

        # Display similar questions if any
        if similar_questions:
            st.markdown("### 💡 您可能还想了解")
            col_count = min(2, len(similar_questions))
            cols = st.columns(col_count)
            for i, question in enumerate(similar_questions):
                with cols[i % col_count]:
                    if st.button(question, key=f"home_similar_{i}"):
                        # When user clicks on a similar question, use it as new input
                        user_input = question
                        st.session_state.messages.append({"role": "user", "content": user_input})
                        with st.chat_message("user"):
                            st.markdown(user_input)

                        # Get assistant response for the similar question
                        with st.chat_message("assistant"):
                            full_response = ""
                            response_placeholder = st.empty()
                            for event in st.session_state.client.chat_stream(
                                user_input,
                                st.session_state.session_id
                            ):
                                content_chunk = event.get("content", "")
                                if content_chunk:
                                    full_response += content_chunk
                                    response_placeholder.markdown(full_response + "▌")

                            # Finalize response
                            response_placeholder.markdown(full_response)
                            st.session_state.messages.append({"role": "assistant", "content": full_response})

                        # Refresh the page to show the new conversation
                        st.rerun()

def render_analysis_page():
    """Render analysis page with file selection and chat"""
    # File selection section
    st.markdown("### 📁 数据文件选择")
    
    # Create csv_files directory if it doesn't exist
    import os
    os.makedirs("csv_files", exist_ok=True)
    
    # Get list of files in csv_files directory
    csv_files_dir = "csv_files"
    available_files = [f for f in os.listdir(csv_files_dir) if f.endswith((".csv", ".xlsx", ".xls"))]
    
    if available_files:
        selected_file = st.selectbox(
            "从以下文件中选择",
            available_files,
            help="选择已处理好的数据文件进行分析",
            key="analysis_file_selector"
        )
        
        if selected_file:
            file_path = os.path.join(csv_files_dir, selected_file)
            st.session_state.uploaded_file = selected_file
            
            # Show file preview
            try:
                if selected_file.endswith(".csv"):
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)

                st.subheader("📊 文件预览")
                st.dataframe(df.head(), width='stretch')
                st.text(f"数据形状：{df.shape[0]} 行 × {df.shape[1]} 列")
                
                # 以表格形式展示所有列名
                st.subheader("📋 列名信息")
                columns_df = pd.DataFrame({
                    '列名': df.columns.tolist(),
                    '数据类型': df.dtypes.tolist()
                })
                st.dataframe(columns_df, width='stretch')
            except Exception as e:
                st.error(f"❌ 读取文件时出错：{str(e)}")
    else:
        st.info("📁 csv_files目录中暂无数据文件，请先在Home界面上传文件")

    # Chat section
    st.markdown("### 💬 对话分析")

    # Display chat history
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Chat input
    user_input = st.chat_input("请输入您的分析需求...", key="analysis_chat_input")

    # Process user input
    if user_input:
        # Check server connection first
        if not st.session_state.client.check_server_status():
            st.error("❌ 无法连接到服务器。请检查服务器URL并确保服务器正在运行。")
            return

        # Add user message to chat
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)

        # Get assistant response
        with st.chat_message("assistant"):
            full_response = ""
            response_placeholder = st.empty()
            
            # 显示正在处理的提示，而不是显示原始内容
            response_placeholder.markdown("⏳ 正在处理您的请求...")
            
            for event in st.session_state.client.chat_stream(
                user_input,
                st.session_state.session_id
            ):
                content_chunk = event.get("content", "")
                if content_chunk:
                    full_response += content_chunk

            # Finalize response - 清空占位符，然后调用处理函数展示
            response_placeholder.empty()
            process_and_display_assistant_message(full_response)
            st.session_state.messages.append({"role": "assistant", "content": full_response})

        # Get similar questions after receiving response
        similar_questions = st.session_state.client.get_similar_questions(
            user_input, 
            st.session_state.messages[-5:]
        )

        # Display similar questions if any
        if similar_questions:
            st.markdown("### 💡 您可能还想了解")
            col_count = min(2, len(similar_questions))
            cols = st.columns(col_count)
            for i, question in enumerate(similar_questions):
                with cols[i % col_count]:
                    if st.button(question, key=f"analysis_similar_{i}"):
                        # When user clicks on a similar question, use it as new input
                        user_input = question
                        st.session_state.messages.append({"role": "user", "content": user_input})
                        with st.chat_message("user"):
                            st.markdown(user_input)

                        # Get assistant response for the similar question
                        with st.chat_message("assistant"):
                            full_response = ""
                            response_placeholder = st.empty()
                            for event in st.session_state.client.chat_stream(
                                user_input,
                                st.session_state.session_id
                            ):
                                content_chunk = event.get("content", "")
                                if content_chunk:
                                    full_response += content_chunk
                                    response_placeholder.markdown(full_response + "▌")

                            # Finalize response
                            response_placeholder.markdown(full_response)
                            st.session_state.messages.append({"role": "assistant", "content": full_response})

                        # Refresh the page to show the new conversation
                        st.rerun()

def render_documents_page():
    """Render my documents page"""
    st.markdown("# 📄 我的文档")
    
    # 1. 心得体会生成功能
    st.markdown("## 📝 心得体会生成")
    with st.form("reflection_form"):
        topic = st.text_input("主题内容（可选）", placeholder="请输入心得体会的主题")
        word_limit = st.number_input("字数限制（可选）", min_value=100, max_value=5000, value=500, step=100)
        outline = st.text_area("文本大纲（可选）", placeholder="请输入文本大纲，每点一行")
        
        submitted = st.form_submit_button("生成心得体会")
        
        if submitted:
            if not topic:
                st.warning("请至少输入主题内容")
            else:
                with st.spinner("正在生成心得体会..."):
                    try:
                        # 构建请求
                        prompt_parts = [f"生成关于{topic}的心得体会"]
                        if word_limit:
                            prompt_parts.append(f"，字数限制{word_limit}字")
                        if outline:
                            prompt_parts.append(f"，提纲包括：{outline}")
                        prompt = "".join(prompt_parts)
                        
                        # 调用API生成
                        full_response = ""
                        response_placeholder = st.empty()
                        for event in st.session_state.client.chat_stream(
                            prompt,
                            st.session_state.session_id
                        ):
                            content_chunk = event.get("content", "")
                            if content_chunk:
                                full_response += content_chunk
                                response_placeholder.markdown(full_response + "▌")
                        
                        response_placeholder.markdown(full_response)
                        
                        # 保存生成的内容
                        os.makedirs("generated_reflections", exist_ok=True)
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')[:20]
                        filename = f"reflection_{safe_topic}_{timestamp}.txt"
                        filepath = os.path.join("generated_reflections", filename)
                        
                        with open(filepath, "w", encoding="utf-8") as f:
                            f.write(full_response)
                        
                        st.success(f"✅ 心得体会已生成并保存到：{filepath}")
                        
                        # 提供下载按钮
                        st.download_button(
                            label="📥 下载心得体会",
                            data=full_response,
                            file_name=filename,
                            mime="text/plain"
                        )
                        
                    except Exception as e:
                        st.error(f"❌ 生成失败：{str(e)}")
    
    st.markdown("---")
    
    # 2. 历史文档查阅功能
    st.markdown("## 📁 历史文档")
    
    # 定义文档目录
    documents_dir = "docx_files"
    os.makedirs(documents_dir, exist_ok=True)
    
    # 搜索功能
    search_query = st.text_input("🔍 搜索文档", placeholder="输入关键词搜索文档...")
    
    # 展示已上传的文件列表
    st.markdown("### 文档列表")
    
    if os.path.exists(documents_dir):
        files = os.listdir(documents_dir)
        
        # 过滤搜索结果
        if search_query:
            files = [f for f in files if search_query.lower() in f.lower()]
        
        if files:
            # 按文件类型分组
            files_by_type = {}
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext not in files_by_type:
                    files_by_type[ext] = []
                files_by_type[ext].append(file)
            
            # 展示每个类型的文件
            for ext, file_list in files_by_type.items():
                st.markdown(f"#### {ext.upper()} 文件 ({len(file_list)})")
                for idx, file in enumerate(file_list, 1):
                    col1, col2, col3 = st.columns([4, 1, 1])
                    with col1:
                        st.markdown(f"**{idx}.** {file}")
                    with col2:
                        if st.button("预览", key=f"preview_{file}"):
                            file_path = os.path.join(documents_dir, file)
                            try:
                                st.subheader(f"📄 文档预览：{file}")
                                
                                # 根据文件类型显示不同的预览
                                ext = os.path.splitext(file)[1].lower()
                                
                                if ext in [".txt", ".md", ".py"]:
                                    # 文本文件预览
                                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                                        content = f.read()
                                    st.code(content, language=ext[1:] if ext == ".py" else "markdown" if ext == ".md" else "plaintext")
                                
                                elif ext == ".csv":
                                    # CSV文件预览
                                    df = pd.read_csv(file_path)
                                    st.dataframe(df, width='stretch')
                                
                                elif ext in [".xlsx", ".xls"]:
                                    # Excel文件预览
                                    df = pd.read_excel(file_path)
                                    st.dataframe(df, width='stretch')
                                
                                elif ext == ".docx":
                                    # Word文档预览
                                    try:
                                        from docx import Document
                                        doc = Document(file_path)
                                        doc_text = "\n".join([para.text for para in doc.paragraphs])
                                        st.markdown(doc_text)
                                    except ImportError:
                                        st.info("请安装python-docx库以预览Word文档")
                                    except Exception as e:
                                        st.error(f"读取Word文档失败：{str(e)}")
                                
                                elif ext in [".ppt", ".pptx"]:
                                    # PowerPoint预览
                                    try:
                                        from pptx import Presentation
                                        prs = Presentation(file_path)
                                        for i, slide in enumerate(prs.slides[:5]):  # 只展示前5页
                                            st.markdown(f"### 幻灯片 {i+1}")
                                            for shape in slide.shapes:
                                                if hasattr(shape, "text"):
                                                    st.markdown(shape.text)
                                    except ImportError:
                                        st.info("请安装python-pptx库以预览PowerPoint文档")
                                    except Exception as e:
                                        st.error(f"读取PowerPoint文档失败：{str(e)}")
                                
                                else:
                                    st.info(f"不支持 {ext} 文件的预览")
                                    
                            except Exception as e:
                                st.error(f"❌ 预览文件失败：{str(e)}")
                    with col3:
                        if st.button("删除", key=f"delete_{file}"):
                            # 删除文件
                            os.remove(os.path.join(documents_dir, file))
                            st.success(f"✅ 已删除：{file}")
                            st.rerun()
        else:
            st.info("📭 暂无文档")
    else:
        st.info("📭 暂无文档")
    
    st.markdown("---")
    
    # 3. 文档上传功能
    st.markdown("## 📤 上传文档")
    uploaded_file = st.file_uploader(
        "选择要上传的文档",
        type=["docx", "txt", "xlsx", "xls", "ppt", "pptx", "md", "py", "csv"],
        help="支持的格式：DOCX, TXT, XLSX, XLS, PPT, PPTX, MD, PY, CSV"
    )
    
    if uploaded_file is not None:
        # 保存文件
        file_path = os.path.join(documents_dir, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.success(f"✅ 文档已上传：{uploaded_file.name}")
        st.rerun()

def render_knowledge_page():
    """Render knowledge base page"""
    st.markdown("# 🧠 知识库")
    
    # 1. 添加网页链接到知识库
    st.markdown("## 🔗 添加网页链接")
    
    with st.form("add_link_form"):
        url = st.text_input("网页链接", placeholder="请输入网页链接（例如：https://example.com）")
        category = st.text_input("分类（可选）", placeholder="请输入分类，如：技术、学习、资讯等")
        submitted = st.form_submit_button("添加到知识库")
        
        if submitted:
            if not url:
                st.error("请输入网页链接")
            else:
                with st.spinner("正在读取网页内容并生成摘要..."):
                    try:
                        import requests
                        from bs4 import BeautifulSoup
                        
                        # 获取网页内容
                        headers = {
                            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                        }
                        response = requests.get(url, headers=headers, timeout=10)
                        response.raise_for_status()
                        
                        # 解析HTML
                        soup = BeautifulSoup(response.text, 'html.parser')
                        
                        # 获取标题
                        title = soup.title.string if soup.title else url
                        
                        # 获取正文内容（提取段落）
                        paragraphs = soup.find_all('p')
                        content_text = '\n'.join([p.get_text().strip() for p in paragraphs[:20] if p.get_text().strip()])
                        
                        # 调用API生成摘要和分类
                        if content_text:
                            prompt = f"""请为以下网页内容生成一个简洁的摘要（约200字），并给出一个合适的分类标签。

                                            网页标题：{title}
                                            网页内容：
                                            {content_text[:3000]}

                                            请按以下格式输出：
                                            【摘要】
                                            （在此处填写摘要）

                                            【分类】
                                            （在此处填写分类标签）
                                            """
                            
                            full_response = ""
                            response_placeholder = st.empty()
                            for event in st.session_state.client.chat_stream(
                                prompt,
                                st.session_state.session_id
                            ):
                                content_chunk = event.get("content", "")
                                if content_chunk:
                                    full_response += content_chunk
                                    response_placeholder.markdown(full_response + "▌")
                            
                            response_placeholder.markdown(full_response)
                            
                            # 使用用户指定的分类或自动生成的分类
                            final_category = category if category else "未分类"
                            
                            # 保存到知识库
                            os.makedirs("knowledge_base", exist_ok=True)
                            file_name = f"link_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
                            file_path = os.path.join("knowledge_base", file_name)
                            
                            with open(file_path, "w", encoding="utf-8") as f:
                                f.write(f"# 📌 网页链接：{url}\n\n")
                                f.write(f"## 📝 标题：{title}\n\n")
                                f.write(f"## 🏷️ 分类：{final_category}\n\n")
                                f.write(f"## 📄 内容摘要：\n\n")
                                f.write(full_response)
                                f.write(f"\n\n---\n\n")
                                f.write(f"## 📅 添加时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                            
                            st.success(f"✅ 链接已成功添加到知识库！")
                            
                            # 显示添加的信息
                            with st.expander("查看保存的内容", expanded=True):
                                st.markdown(f"**链接：** {url}")
                                st.markdown(f"**标题：** {title}")
                                st.markdown(f"**分类：** {final_category}")
                                st.markdown(full_response)
                        else:
                            st.warning("未能从网页中提取到有效内容，请检查链接是否正确")
                    
                    except requests.exceptions.RequestException as e:
                        st.error(f"❌ 获取网页失败：{str(e)}")
                    except Exception as e:
                        st.error(f"❌ 处理失败：{str(e)}")
    
    st.markdown("---")
    
    # 2. 知识库内容展示和管理
    st.markdown("## 📚 知识库内容")
    
    # 搜索和筛选
    col1, col2 = st.columns([3, 1])
    with col1:
        search_kb = st.text_input("🔍 搜索知识库", placeholder="输入关键词搜索...")
    with col2:
        # 获取所有分类
        categories = ["全部"]
        if os.path.exists("knowledge_base"):
            files = os.listdir("knowledge_base")
            for file in files:
                if file.endswith(".md"):
                    try:
                        with open(os.path.join("knowledge_base", file), 'r', encoding='utf-8') as f:
                            content = f.read()
                            category_line = next((line for line in content.split('\n') if line.startswith('## 🏷️ 分类：')), '')
                            if category_line:
                                cat = category_line.replace('## 🏷️ 分类：', '').strip()
                                if cat not in categories:
                                    categories.append(cat)
                    except:
                        pass
        
        selected_category = st.selectbox("📂 分类筛选", categories)
    
    # 展示知识库内容
    if os.path.exists("knowledge_base"):
        files = [f for f in os.listdir("knowledge_base") if f.endswith(".md")]
        
        if files:
            # 过滤文件
            filtered_files = []
            for file in files:
                try:
                    with open(os.path.join("knowledge_base", file), 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    # 搜索过滤
                    if search_kb and search_kb.lower() not in content.lower():
                        continue
                    
                    # 分类过滤
                    if selected_category != "全部":
                        category_line = next((line for line in content.split('\n') if line.startswith('## 🏷️ 分类：')), '')
                        file_category = category_line.replace('## 🏷️ 分类：', '').strip() if category_line else '未分类'
                        if file_category != selected_category:
                            continue
                    
                    filtered_files.append((file, content))
                except Exception as e:
                    continue
            
            if filtered_files:
                st.markdown(f"### 📊 共找到 {len(filtered_files)} 条记录")
                
                # 展示每个知识库条目
                for idx, (file, content) in enumerate(filtered_files, 1):
                    try:
                        lines = content.split('\n')
                        
                        # 提取信息
                        url_line = next((line for line in lines if line.startswith('# 📌 网页链接：')), '')
                        title_line = next((line for line in lines if line.startswith('## 📝 标题：')), '')
                        category_line = next((line for line in lines if line.startswith('## 🏷️ 分类：')), '')
                        time_line = next((line for line in lines if line.startswith('## 📅 添加时间：')), '')
                        
                        url = url_line.replace('# 📌 网页链接：', '').strip() if url_line else '未知链接'
                        title = title_line.replace('## 📝 标题：', '').strip() if title_line else '无标题'
                        category = category_line.replace('## 🏷️ 分类：', '').strip() if category_line else '未分类'
                        add_time = time_line.replace('## 📅 添加时间：', '').strip() if time_line else ''
                        
                        # 创建expander
                        expander_title = f"📑 {title[:50]}..." if len(title) > 50 else f"📑 {title}"
                        with st.expander(expander_title, expanded=False):
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                st.markdown(f"**🔗 链接：** [{url}]({url})")
                                st.markdown(f"**🏷️ 分类：** {category}")
                                if add_time:
                                    st.markdown(f"**📅 添加时间：** {add_time}")
                            with col2:
                                if st.button(f"🗑️ 删除", key=f"delete_kb_{idx}"):
                                    os.remove(os.path.join("knowledge_base", file))
                                    st.success(f"✅ 已删除：{title}")
                                    st.rerun()
                            
                            st.markdown("---")
                            st.markdown("**📄 内容摘要：**")
                            st.markdown(content)
                    
                    except Exception as e:
                        st.error(f"❌ 读取文件失败：{str(e)}")
            else:
                st.info("� 没有找到匹配的知识库条目")
        else:
            st.info("�� 知识库为空，快去添加一些网页链接吧！")
    else:
        st.info("📭 知识库为空，快去添加一些网页链接吧！")
    
    st.markdown("---")
    
    # 3. 常用分析模板
    st.markdown("## 📋 常用分析模板")
    templates = [
        {"name": "时限分析", "desc": "分析各省份、各产品类型时限"},
        {"name": "物流时效分析", "desc": "分析不同线路、不同时段的物流时效"},
    ]

    for template in templates:
        with st.expander(f"{template['name']}", expanded=False):
            st.markdown(f"**描述：** {template['desc']}")
            if st.button(f"使用模板：{template['name']}", key=f"template_{template['name']}"):
                # 自动填充聊天框
                st.session_state.messages.append({"role": "user", "content": f"按照{template['name']}模板分析我的数据"})
                st.rerun()

def render_python_sandbox_page():
    """Render Python sandbox page"""
    st.markdown("# 🐍 Python 代码沙盒")
    st.markdown("---")
    
    # 代码编辑器区域
    st.markdown("## 📝 代码编辑器")
    
    # 初始化代码
    if "sandbox_code" not in st.session_state:
        st.session_state.sandbox_code = """import pandas as pd
        import numpy as np

        # 示例代码
        print("欢迎使用Python沙盒！")

        # 创建示例数据
        data = {
            '姓名': ['张三', '李四', '王五', '赵六'],
            '年龄': [25, 30, 35, 28],
            '城市': ['北京', '上海', '广州', '深圳']
        }
        df = pd.DataFrame(data)
        print("\\n示例数据:")
        print(df)
        print("\\n统计信息:")
        print(df.describe())
        """
    
    code = st.text_area(
        "输入Python代码",
        value=st.session_state.sandbox_code,
        height=300,
        key="python_sandbox_editor"
    )
    st.session_state.sandbox_code = code
    
    # 执行按钮
    col1, col2 = st.columns([1, 4])
    with col1:
        execute_btn = st.button("▶️ 执行代码", type="primary", key="execute_sandbox_code")
    
    with col2:
        clear_btn = st.button("🗑️ 清空代码", key="clear_sandbox_code")
        if clear_btn:
            st.session_state.sandbox_code = ""
            st.rerun()
    
    # 执行代码
    if execute_btn:
        if code.strip():
            with st.spinner("正在执行代码..."):
                try:
                    # 使用code_manager执行
                    from src.utils.code_manager import execute_python_sandbox
                    
                    # 设置工作目录
                    workspace_dir = os.getcwd()
                    result = execute_python_sandbox(code, workspace_dir)
                    
                    st.markdown("---")
                    if result["success"]:
                        st.success("✅ 执行成功！")
                        st.markdown(f"**执行时间:** {result['execution_time']:.2f} 秒")
                        
                        if result["output"]:
                            st.markdown("### 📤 输出结果")
                            st.code(result["output"], language="text")
                    else:
                        st.error("❌ 执行失败")
                        if result["error"]:
                            st.markdown("### 🚨 错误信息")
                            st.code(result["error"], language="text")
                            
                except Exception as e:
                    st.error(f"执行错误: {str(e)}")
                    st.exception(e)
        else:
            st.warning("请输入Python代码")
    
    st.markdown("---")
    
    # 已保存的代码文件区域
    st.markdown("## 📁 已保存的代码文件")
    
    # 查找工作区中的generated_scripts目录
    generated_scripts_dir = "generated_scripts"
    if os.path.exists(generated_scripts_dir):
        code_files = [f for f in os.listdir(generated_scripts_dir) if f.endswith('.py')]
        
        if code_files:
            # 按修改时间排序
            code_files.sort(key=lambda x: os.path.getmtime(os.path.join(generated_scripts_dir, x)), reverse=True)
            
            for filename in code_files:
                file_path = os.path.join(generated_scripts_dir, filename)
                with st.expander(f"📄 {filename}", expanded=False):
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            file_content = f.read()
                        
                        st.code(file_content, language='python', line_numbers=True)
                        
                        col1, col2, col3 = st.columns([1, 1, 2])
                        with col1:
                            if st.button(f"📥 加载", key=f"load_{filename}"):
                                st.session_state.sandbox_code = file_content
                                st.rerun()
                        with col2:
                            if st.button(f"🗑️ 删除", key=f"del_{filename}"):
                                os.remove(file_path)
                                st.success(f"已删除: {filename}")
                                st.rerun()
                        with col3:
                            st.text(f"修改时间: {datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%Y-%m-%d %H:%M:%S')}")
                            
                    except Exception as e:
                        st.error(f"读取文件失败: {str(e)}")
        else:
            st.info("📭 暂无保存的代码文件")
    else:
        st.info("📭 generated_scripts目录不存在，请先运行一些分析任务")
    
    st.markdown("---")
    
    # 常用代码模板
    st.markdown("## 📋 常用代码模板")
    
    templates = [
        {
            "name": "数据加载和预览",
            "code": """import pandas as pd

            # 读取CSV文件
            df = pd.read_csv('csv_files/您的文件名.csv')
            print("数据形状:", df.shape)
            print("\\n前5行数据:")
            print(df.head())
            print("\\n数据类型:")
            print(df.dtypes)
            print("\\n统计信息:")
            print(df.describe())"""
                    },
                    {
                        "name": "数据可视化",
                        "code": """import pandas as pd
            import matplotlib.pyplot as plt
            import seaborn as sns

            # 读取数据
            df = pd.read_csv('csv_files/您的文件名.csv')

            # 设置中文显示
            plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei']
            plt.rcParams['axes.unicode_minus'] = False

            # 创建简单图表
            plt.figure(figsize=(10, 6))
            # 替换'列名'为实际的列名
            if '列名' in df.columns:
                sns.histplot(df['列名'], bins=30, kde=True)
                plt.title('数据分布图')
                plt.xlabel('值')
                plt.ylabel('频数')
                plt.tight_layout()
                plt.show()
            else:
                print("请替换'列名'为实际的列名")
                print("可用列:", df.columns.tolist())"""
                    },
                    {
                        "name": "数据清洗",
                        "code": """import pandas as pd
            import numpy as np

            # 读取数据
            df = pd.read_csv('csv_files/您的文件名.csv')

            print("原始数据形状:", df.shape)
            print("\\n缺失值统计:")
            print(df.isnull().sum())

            # 处理缺失值
            # 方法1: 删除含有缺失值的行
            df_clean = df.dropna()
            print("\\n删除缺失值后形状:", df_clean.shape)

            # 方法2: 用均值填充数值列
            # for col in df.select_dtypes(include=[np.number]).columns:
            #     df[col] = df[col].fillna(df[col].mean())

            # 方法3: 用众数填充分类列
            # for col in df.select_dtypes(include=['object']).columns:
            #     df[col] = df[col].fillna(df[col].mode()[0])

            print("\\n清洗完成！")"""
        }
    ]
    
    for i, template in enumerate(templates):
        with st.expander(f"📑 {template['name']}", expanded=False):
            st.code(template['code'], language='python')
            if st.button(f"使用此模板", key=f"template_{i}"):
                st.session_state.sandbox_code = template['code']
                st.rerun()

def render_history_page():
    """Render chat history page"""
    st.markdown("# 📚 历史会话")
    if st.session_state.messages:
        # 按时间倒序展示
        for idx, message in enumerate(reversed(st.session_state.messages)):
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
            # 每两条消息分隔
            if idx % 2 == 1 and idx != len(st.session_state.messages)-1:
                st.markdown("---")
    else:
        st.info("📭 暂无历史会话")

def main():
    """Main application entry point"""
    # Initialize session state
    initialize_session_state()

    # Render sidebar
    render_sidebar()

    # Render main content based on current page
    current_page = st.session_state.current_page

    if current_page == "home":
        render_home_page()
    elif current_page == "analysis":
        render_analysis_page()
    elif current_page == "python_sandbox":
        render_python_sandbox_page()
    elif current_page == "documents":
        render_documents_page()
    elif current_page == "knowledge":
        render_knowledge_page()
    elif current_page == "history":
        render_history_page()

if __name__ == "__main__":
    main()