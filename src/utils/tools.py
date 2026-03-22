import json
import logging
import os
import platform
from pathlib import Path
from typing import Optional, Literal

from langchain_core.tools import tool

from src.config.loader import load_yaml_config

logger = logging.getLogger(__name__)

# Limit the data size for single read operation (in bytes), can be overridden by environment variable READ_FILE_MAX_BYTES
_DEFAULT_MAX_BYTES = 32768  # 32 KB
try:
    _MAX_RETURN_BYTES = int(os.environ.get('READ_FILE_MAX_BYTES', str(_DEFAULT_MAX_BYTES)))
except Exception:
    _MAX_RETURN_BYTES = _DEFAULT_MAX_BYTES


def _ensure_within_limit_by_bytes(text: str, *, context: str) -> None:
    """Raise an exception if the text exceeds the limit when encoded in UTF-8."""
    if text is None:
        return
    if len(text.encode('utf-8', errors='ignore')) > _MAX_RETURN_BYTES:
        raise ValueError(f"Read result is too large, exceeding token limit")


@tool
def ask_user(question: str):
    """Ask the user for information needed to continue analysis.

    Args:
        question (str): Clear question describing what information is needed and why
    """
    return


@tool
def feedback(issue: str):
    """Report issues when required data cannot be obtained or conditions cannot be met during analysis execution

    Args:
        issue (str): A clear description of the encountered issue or obstacle
    """
    return


@tool
def terminate(status: str):
    """When you have finished the step, call this tool to end the work.

    Args:
        status (str): The finish status of the step. ["success", "failure"]
    """
    content = f"The step has been completed with status: {str(status)}"
    logger.info(content)
    return content


@tool
def read_file_head3(file_path: str):
    """
    Read the first 3 lines of a file

    Args:
        file_path (str): File path

    Returns:
        str: First 3 lines of the file
    """
    return read_file(file_path, "head", 3)


@tool
def read_file_head20(file_path: str):
    """
    Read the first 20 lines of a file

    Args:
        file_path (str): File path

    Returns:
        str: First 20 lines of the file
    """
    return read_file(file_path, "head", 20)


@tool
def list_available_data_files() -> str:
    """
    List all available data files (CSV and Excel).

    Returns:
        str: JSON string containing array of data file information
    """
    try:
        # Load configuration
        config = load_yaml_config("conf.yaml")
        app_conf = config.get("app", {}) or {}
        csv_conf = app_conf.get("csv_data_directory", {}) or {}
        
        # Get data directory based on OS
        system = platform.system()
        if system == "Windows":
            key = "windows"
        elif system == "Darwin":
            key = "macOS"
        else:
            key = "linux"
        base_dir = csv_conf.get(key)
        
        if not base_dir:
            return json.dumps({
                "error": "Data directory not configured in conf.yaml",
                "files": []
            }, ensure_ascii=False, indent=2)
        
        if not os.path.isdir(base_dir):
            return json.dumps({
                "error": f"Data directory does not exist: {base_dir}",
                "files": []
            }, ensure_ascii=False, indent=2)
        
        # Scan directory for data files
        path = Path(base_dir)
        data_files = list(path.glob('*.csv')) + list(path.glob('*.xlsx')) + list(path.glob('*.xls'))
        
        files_info = []
        for file_path in data_files:
            try:
                # Get basic file info
                file_info = {
                    "file_name": file_path.name,
                    "file_path": str(file_path).replace('\\', '/'),
                    "file_type": file_path.suffix.lower(),
                    "columns": []
                }
                
                # Try to read file to get column info
                try:
                    import pandas as pd
                    if file_path.suffix.lower() in ['.csv']:
                        # Try different separators for CSV
                        separators = ['^', ',', '\t', ';', '|']
                        df_header = None
                        for sep in separators:
                            try:
                                df_header = pd.read_csv(file_path, sep=sep, nrows=0)  # Only read header
                                if len(df_header.columns) > 1:
                                    break
                            except:
                                continue
                    elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                        # Read Excel file
                        df_header = pd.read_excel(file_path, nrows=0)
                    else:
                        continue
                    
                    if df_header is not None and len(df_header.columns) > 1:
                        file_info["columns"] = [{"name": col, "dtype": str(df_header[col].dtype)} for col in df_header.columns]
                except Exception as e:
                    logger.warning(f"Failed to analyze file {file_path}: {e}")
                    # Still include file info even if analysis fails
                
                files_info.append(file_info)
            except Exception as e:
                logger.warning(f"Failed to process file {file_path}: {e}")
                continue
        
        return json.dumps({
            "directory": str(base_dir).replace('\\', '/'),
            "files": files_info
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Failed to list data files: {e}", exc_info=True)
        return json.dumps({
            "error": f"Failed to list data files: {str(e)}",
            "files": []
        }, ensure_ascii=False, indent=2)


@tool
def read_data_file(file_path: str, n_rows: int = None, sheet_name: str = None) -> str:
    """
    Read data file (CSV or Excel) with auto-detection of separator.
    IMPORTANT: By default, reads ALL rows from the file (full dataset).
    Only specify n_rows if you need a limited sample for testing.
    
    If file_path is just a filename (without directory), will search in:
    1. Current directory
    2. csv_files/ directory (for uploaded files)

    Args:
        file_path (str): Path to data file or filename
        n_rows (int): Number of rows to read. DEFAULT: None (reads ALL rows). Set to specific number only for sampling
        sheet_name (str): Sheet name for Excel files (default: first sheet)

    Returns:
        str: JSON string with data and metadata including row_count
    """
    try:
        import pandas as pd
        
        # Convert to Path object for easier manipulation
        path_obj = Path(file_path)
        
        # Check if file exists at the given path
        if not path_obj.exists():
            # If it's just a filename (no directory components), search in csv_files/
            if path_obj.parent == Path('.'):
                csv_files_dir = Path("csv_files") / path_obj
                if csv_files_dir.exists():
                    file_path = str(csv_files_dir)
                    path_obj = csv_files_dir
                    logger.info(f"Found file in csv_files directory: {file_path}")
                else:
                    # Also check with absolute path from project root
                    import sys
                    project_root = Path(__file__).parent.parent.parent
                    abs_path = project_root / "csv_files" / path_obj
                    if abs_path.exists():
                        file_path = str(abs_path)
                        path_obj = abs_path
                        logger.info(f"Found file with absolute path: {file_path}")
                    else:
                        return json.dumps({
                            "error": f"File not found: {file_path}",
                            "hint": f"File '{path_obj.name}' was not found. Please check if the file exists in the csv_files/ directory or provide the correct path."
                        }, ensure_ascii=False)
            else:
                return json.dumps({
                    "error": f"File not found: {file_path}",
                    "hint": f"Please verify the file path is correct and the file exists."
                }, ensure_ascii=False)
        
        df = None
        separator = None
        file_type = Path(file_path).suffix.lower()
        
        if file_type == '.csv':
            # Try different separators for CSV
            separators = ['^', ',', '\t', ';', '|']
            for sep in separators:
                try:
                    df = pd.read_csv(file_path, sep=sep, nrows=n_rows)
                    if len(df.columns) > 1:
                        separator = sep
                        break
                except:
                    continue
            
            if df is None:
                return json.dumps({"error": "Failed to read CSV file with any separator"}, ensure_ascii=False)
        elif file_type in ['.xlsx', '.xls']:
            # Read Excel file
            try:
                df = pd.read_excel(file_path, nrows=n_rows, sheet_name=sheet_name)
            except Exception as e:
                return json.dumps({"error": f"Failed to read Excel file: {str(e)}"}, ensure_ascii=False)
        else:
            return json.dumps({"error": "Unsupported file type. Only CSV and Excel files are supported"}, ensure_ascii=False)
        
        # Clean column names (remove spaces, special characters)
        df.columns = df.columns.str.strip()
        
        # Convert to JSON
        data = {
            "file_path": file_path,
            "file_type": file_type,
            "separator": separator,
            "columns": list(df.columns),
            "row_count": len(df),
            "sample_data": df.head(min(10, len(df))).to_dict('records')
        }
        
        return json.dumps(data, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error reading data file: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def filter_data_file(file_path: str, filter_conditions: str, n_rows: int = 100, sheet_name: str = None) -> str:
    """
    Filter data file (CSV or Excel) based on conditions

    Args:
        file_path (str): Path to data file
        filter_conditions (str): Filter conditions as Python expression (e.g., "column1 > 10 and column2 == 'A'")
        n_rows (int): Maximum number of rows to return (default: 100)
        sheet_name (str): Sheet name for Excel files (default: first sheet)

    Returns:
        str: JSON string with filtered data
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        df = None
        file_type = Path(file_path).suffix.lower()
        
        if file_type == '.csv':
            # Try different separators for CSV
            separators = ['^', ',', '\t', ';', '|']
            for sep in separators:
                try:
                    df = pd.read_csv(file_path, sep=sep)
                    if len(df.columns) > 1:
                        break
                except:
                    continue
            
            if df is None:
                return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        elif file_type in ['.xlsx', '.xls']:
            # Read Excel file
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
            except Exception as e:
                return json.dumps({"error": f"Failed to read Excel file: {str(e)}"}, ensure_ascii=False)
        else:
            return json.dumps({"error": "Unsupported file type. Only CSV and Excel files are supported"}, ensure_ascii=False)
        
        # Clean column names
        df.columns = df.columns.str.strip()
        
        # Apply filter
        try:
            filtered_df = df.query(filter_conditions)
        except Exception as e:
            return json.dumps({"error": f"Filter condition error: {str(e)}"}, ensure_ascii=False)
        
        # Convert to JSON
        data = {
            "file_path": file_path,
            "file_type": file_type,
            "filter_conditions": filter_conditions,
            "total_rows": len(filtered_df),
            "data": filtered_df.head(n_rows).to_dict('records')
        }
        
        return json.dumps(data, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error filtering data file: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def analyze_data_statistics(file_path: str, columns: str, sheet_name: str = None) -> str:
    """
    Analyze statistics for specified columns in data file (CSV or Excel)

    Args:
        file_path (str): Path to data file
        columns (str): Comma-separated list of columns to analyze
        sheet_name (str): Sheet name for Excel files (default: first sheet)

    Returns:
        str: JSON string with statistics
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        df = None
        file_type = Path(file_path).suffix.lower()
        
        if file_type == '.csv':
            # Try different separators for CSV
            separators = ['^', ',', '\t', ';', '|']
            for sep in separators:
                try:
                    df = pd.read_csv(file_path, sep=sep)
                    if len(df.columns) > 1:
                        break
                except:
                    continue
            
            if df is None:
                return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        elif file_type in ['.xlsx', '.xls']:
            # Read Excel file
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
            except Exception as e:
                return json.dumps({"error": f"Failed to read Excel file: {str(e)}"}, ensure_ascii=False)
        else:
            return json.dumps({"error": "Unsupported file type. Only CSV and Excel files are supported"}, ensure_ascii=False)
        
        # Clean column names
        df.columns = df.columns.str.strip()
        
        # Get columns to analyze
        column_list = [col.strip() for col in columns.split(',')]
        
        # Calculate statistics
        stats = {}
        for col in column_list:
            if col in df.columns:
                col_stats = {
                    "count": int(df[col].count()),
                    "mean": float(df[col].mean()) if pd.api.types.is_numeric_dtype(df[col]) else "N/A",
                    "median": float(df[col].median()) if pd.api.types.is_numeric_dtype(df[col]) else "N/A",
                    "std": float(df[col].std()) if pd.api.types.is_numeric_dtype(df[col]) else "N/A",
                    "min": float(df[col].min()) if pd.api.types.is_numeric_dtype(df[col]) else str(df[col].min()),
                    "max": float(df[col].max()) if pd.api.types.is_numeric_dtype(df[col]) else str(df[col].max()),
                    "unique": int(df[col].nunique()),
                    "missing": int(df[col].isnull().sum())
                }
                # Add quartiles for numeric columns
                if pd.api.types.is_numeric_dtype(df[col]):
                    col_stats["q25"] = float(df[col].quantile(0.25))
                    col_stats["q75"] = float(df[col].quantile(0.75))
                stats[col] = col_stats
            else:
                stats[col] = {"error": "Column not found"}
        
        data = {
            "file_path": file_path,
            "file_type": file_type,
            "columns": column_list,
            "statistics": stats
        }
        
        return json.dumps(data, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error analyzing data statistics: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def calculate_logistics_metrics(file_path: str, actual_time_col: str, planned_time_col: str) -> str:
    """
    Calculate logistics performance metrics

    Args:
        file_path (str): Path to CSV file
        actual_time_col (str): Column name for actual delivery time
        planned_time_col (str): Column name for planned delivery time

    Returns:
        str: JSON string with logistics metrics
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        if actual_time_col not in df.columns:
            return json.dumps({"error": f"Column not found: {actual_time_col}"}, ensure_ascii=False)
        if planned_time_col not in df.columns:
            return json.dumps({"error": f"Column not found: {planned_time_col}"}, ensure_ascii=False)
        
        # Calculate metrics
        total_orders = len(df)
        
        # Calculate on-time rate
        df['is_on_time'] = df[actual_time_col] <= df[planned_time_col]
        on_time_count = df['is_on_time'].sum()
        on_time_rate = (on_time_count / total_orders) * 100 if total_orders > 0 else 0
        
        # Calculate delay rate
        delay_rate = 100 - on_time_rate
        
        # Calculate average delay time for delayed orders
        df['delay_time'] = df[actual_time_col] - df[planned_time_col]
        delayed_orders = df[df['delay_time'] > 0]
        avg_delay_time = delayed_orders['delay_time'].mean() if len(delayed_orders) > 0 else 0
        
        # Calculate delivery time statistics
        delivery_time_stats = {
            "mean": float(df[actual_time_col].mean()),
            "median": float(df[actual_time_col].median()),
            "std": float(df[actual_time_col].std()),
            "min": float(df[actual_time_col].min()),
            "max": float(df[actual_time_col].max())
        }
        
        metrics = {
            "total_orders": total_orders,
            "on_time_count": int(on_time_count),
            "on_time_rate": round(on_time_rate, 2),
            "delay_rate": round(delay_rate, 2),
            "delayed_orders": len(delayed_orders),
            "average_delay_time": round(float(avg_delay_time), 2) if avg_delay_time else 0,
            "delivery_time_statistics": delivery_time_stats
        }
        
        return json.dumps(metrics, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error calculating logistics metrics: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def analyze_trend_data(file_path: str, date_col: str, value_col: str, period: str = "daily") -> str:
    """
    Analyze trend data over time

    Args:
        file_path (str): Path to CSV file
        date_col (str): Column name for date/time
        value_col (str): Column name for value to analyze
        period (str): Time period for analysis (daily, weekly, monthly)

    Returns:
        str: JSON string with trend analysis
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        if date_col not in df.columns:
            return json.dumps({"error": f"Column not found: {date_col}"}, ensure_ascii=False)
        if value_col not in df.columns:
            return json.dumps({"error": f"Column not found: {value_col}"}, ensure_ascii=False)
        
        # Convert date column to datetime
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        
        # Group by period
        if period == "daily":
            df['period'] = df[date_col].dt.date
        elif period == "weekly":
            df['period'] = df[date_col].dt.to_period('W')
        elif period == "monthly":
            df['period'] = df[date_col].dt.to_period('M')
        else:
            return json.dumps({"error": "Invalid period. Use daily, weekly, or monthly"}, ensure_ascii=False)
        
        # Calculate trend statistics
        trend_data = df.groupby('period')[value_col].agg(['mean', 'median', 'std', 'count']).reset_index()
        trend_data['period'] = trend_data['period'].astype(str)
        
        # Calculate growth rate
        if len(trend_data) > 1:
            trend_data['growth_rate'] = trend_data['mean'].pct_change() * 100
        else:
            trend_data['growth_rate'] = 0
        
        result = {
            "period": period,
            "trend_data": trend_data.to_dict('records'),
            "overall_stats": {
                "mean": float(df[value_col].mean()),
                "median": float(df[value_col].median()),
                "std": float(df[value_col].std()),
                "total_count": len(df)
            }
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error analyzing trend data: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def calculate_route_performance(file_path: str, route_col: str, actual_time_col: str, planned_time_col: str) -> str:
    """
    Calculate performance metrics by route

    Args:
        file_path (str): Path to CSV file
        route_col (str): Column name for route identifier
        actual_time_col (str): Column name for actual delivery time
        planned_time_col (str): Column name for planned delivery time

    Returns:
        str: JSON string with route performance metrics
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        for col in [route_col, actual_time_col, planned_time_col]:
            if col not in df.columns:
                return json.dumps({"error": f"Column not found: {col}"}, ensure_ascii=False)
        
        # Calculate route performance
        df['is_on_time'] = df[actual_time_col] <= df[planned_time_col]
        df['delay_time'] = df[actual_time_col] - df[planned_time_col]
        
        route_performance = df.groupby(route_col).agg({
            'is_on_time': ['sum', 'count'],
            'delay_time': 'mean',
            actual_time_col: ['mean', 'std']
        }).reset_index()
        
        # Flatten multi-index columns
        route_performance.columns = ['route', 'on_time_count', 'total_count', 'avg_delay_time', 'avg_delivery_time', 'delivery_time_std']
        
        # Calculate on-time rate
        route_performance['on_time_rate'] = (route_performance['on_time_count'] / route_performance['total_count']) * 100
        
        # Sort by on-time rate
        route_performance = route_performance.sort_values('on_time_rate', ascending=False)
        
        result = {
            "route_performance": route_performance.to_dict('records'),
            "summary": {
                "total_routes": len(route_performance),
                "average_on_time_rate": float(route_performance['on_time_rate'].mean()),
                "best_route": route_performance.iloc[0]['route'] if len(route_performance) > 0 else "N/A",
                "worst_route": route_performance.iloc[-1]['route'] if len(route_performance) > 0 else "N/A"
            }
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error calculating route performance: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def calculate_yoy_mom_on_time_rate(file_path: str, date_col: str, actual_time_col: str, planned_time_col: str) -> str:
    """
    Calculate year-over-year (YoY) and month-over-month (MoM) on-time rates

    Args:
        file_path (str): Path to CSV file
        date_col (str): Column name for date
        actual_time_col (str): Column name for actual delivery time
        planned_time_col (str): Column name for planned delivery time

    Returns:
        str: JSON string with YoY and MoM on-time rates
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        for col in [date_col, actual_time_col, planned_time_col]:
            if col not in df.columns:
                return json.dumps({"error": f"Column not found: {col}"}, ensure_ascii=False)
        
        # Convert date column to datetime
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        
        # Remove rows with invalid dates
        df = df.dropna(subset=[date_col])
        
        # Calculate on-time status
        df['is_on_time'] = df[actual_time_col] <= df[planned_time_col]
        
        # Extract year and month
        df['year'] = df[date_col].dt.year
        df['month'] = df[date_col].dt.month
        
        # Calculate monthly on-time rates
        monthly_rates = df.groupby(['year', 'month'])['is_on_time'].agg(['sum', 'count']).reset_index()
        monthly_rates['on_time_rate'] = (monthly_rates['sum'] / monthly_rates['count']) * 100
        
        # Calculate YoY and MoM changes
        monthly_rates = monthly_rates.sort_values(['year', 'month'])
        monthly_rates['mom_change'] = monthly_rates['on_time_rate'].pct_change() * 100
        
        # Calculate YoY change
        result = {
            "monthly_on_time_rates": [],
            "yoy_comparisons": []
        }
        
        # Generate monthly data
        for _, row in monthly_rates.iterrows():
            result["monthly_on_time_rates"].append({
                "year": int(row['year']),
                "month": int(row['month']),
                "on_time_rate": round(float(row['on_time_rate']), 2),
                "mom_change": round(float(row['mom_change']), 2) if not pd.isna(row['mom_change']) else None
            })
        
        # Calculate YoY comparisons
        years = sorted(monthly_rates['year'].unique())
        if len(years) > 1:
            for month in range(1, 13):
                for year in years[1:]:
                    current_year_data = monthly_rates[(monthly_rates['year'] == year) & (monthly_rates['month'] == month)]
                    previous_year_data = monthly_rates[(monthly_rates['year'] == year - 1) & (monthly_rates['month'] == month)]
                    
                    if not current_year_data.empty and not previous_year_data.empty:
                        current_rate = current_year_data.iloc[0]['on_time_rate']
                        previous_rate = previous_year_data.iloc[0]['on_time_rate']
                        yoy_change = ((current_rate - previous_rate) / previous_rate) * 100
                        
                        result["yoy_comparisons"].append({
                            "month": month,
                            "current_year": year,
                            "previous_year": year - 1,
                            "current_rate": round(float(current_rate), 2),
                            "previous_rate": round(float(previous_rate), 2),
                            "yoy_change": round(float(yoy_change), 2)
                        })
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error calculating YoY/MoM on-time rates: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def calculate_time_limit_compliance(file_path: str, actual_time_col: str, planned_time_col: str) -> str:
    """
    Calculate time limit compliance metrics

    Args:
        file_path (str): Path to CSV file
        actual_time_col (str): Column name for actual delivery time
        planned_time_col (str): Column name for planned delivery time

    Returns:
        str: JSON string with time limit compliance metrics
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        for col in [actual_time_col, planned_time_col]:
            if col not in df.columns:
                return json.dumps({"error": f"Column not found: {col}"}, ensure_ascii=False)
        
        # Handle missing values
        df = df.dropna(subset=[actual_time_col, planned_time_col])
        
        # Convert to numeric if needed
        df[actual_time_col] = pd.to_numeric(df[actual_time_col], errors='coerce')
        df[planned_time_col] = pd.to_numeric(df[planned_time_col], errors='coerce')
        
        # Remove rows with invalid numeric values
        df = df.dropna(subset=[actual_time_col, planned_time_col])
        
        # Calculate compliance metrics
        total_orders = len(df)
        
        # On-time (actual <= planned)
        on_time = df[df[actual_time_col] <= df[planned_time_col]]
        on_time_rate = (len(on_time) / total_orders) * 100 if total_orders > 0 else 0
        
        # Slight delay (actual <= planned * 1.1)
        slight_delay = df[(df[actual_time_col] > df[planned_time_col]) & (df[actual_time_col] <= df[planned_time_col] * 1.1)]
        slight_delay_rate = (len(slight_delay) / total_orders) * 100 if total_orders > 0 else 0
        
        # Significant delay (actual > planned * 1.1)
        significant_delay = df[df[actual_time_col] > df[planned_time_col] * 1.1]
        significant_delay_rate = (len(significant_delay) / total_orders) * 100 if total_orders > 0 else 0
        
        # Calculate average delay times
        df['delay_time'] = df[actual_time_col] - df[planned_time_col]
        avg_delay = df[df['delay_time'] > 0]['delay_time'].mean() if len(df[df['delay_time'] > 0]) > 0 else 0
        
        result = {
            "total_orders": total_orders,
            "on_time_rate": round(on_time_rate, 2),
            "slight_delay_rate": round(slight_delay_rate, 2),
            "significant_delay_rate": round(significant_delay_rate, 2),
            "average_delay_time": round(float(avg_delay), 2) if avg_delay else 0,
            "compliance_breakdown": [
                {"category": "On Time", "rate": round(on_time_rate, 2), "count": len(on_time)},
                {"category": "Slight Delay (<=10%)", "rate": round(slight_delay_rate, 2), "count": len(slight_delay)},
                {"category": "Significant Delay (>10%)", "rate": round(significant_delay_rate, 2), "count": len(significant_delay)}
            ]
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error calculating time limit compliance: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def analyze_key_routes(file_path: str, route_col: str, actual_time_col: str, planned_time_col: str, top_n: int = 10) -> str:
    """
    Analyze key routes based on performance metrics

    Args:
        file_path (str): Path to CSV file
        route_col (str): Column name for route identifier
        actual_time_col (str): Column name for actual delivery time
        planned_time_col (str): Column name for planned delivery time
        top_n (int): Number of top routes to analyze

    Returns:
        str: JSON string with key route analysis
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        for col in [route_col, actual_time_col, planned_time_col]:
            if col not in df.columns:
                return json.dumps({"error": f"Column not found: {col}"}, ensure_ascii=False)
        
        # Handle missing values
        df = df.dropna(subset=[route_col, actual_time_col, planned_time_col])
        
        # Convert to numeric if needed
        df[actual_time_col] = pd.to_numeric(df[actual_time_col], errors='coerce')
        df[planned_time_col] = pd.to_numeric(df[planned_time_col], errors='coerce')
        
        # Remove rows with invalid numeric values
        df = df.dropna(subset=[actual_time_col, planned_time_col])
        
        # Calculate route performance
        df['is_on_time'] = df[actual_time_col] <= df[planned_time_col]
        df['delay_time'] = df[actual_time_col] - df[planned_time_col]
        
        route_metrics = df.groupby(route_col).agg({
            'is_on_time': ['sum', 'count'],
            'delay_time': 'mean',
            actual_time_col: ['mean', 'std'],
            planned_time_col: 'mean'
        }).reset_index()
        
        # Flatten columns
        route_metrics.columns = ['route', 'on_time_count', 'total_count', 'avg_delay_time', 'avg_actual_time', 'actual_time_std', 'avg_planned_time']
        
        # Calculate metrics
        route_metrics['on_time_rate'] = (route_metrics['on_time_count'] / route_metrics['total_count']) * 100
        route_metrics['time_efficiency'] = route_metrics['avg_actual_time'] / route_metrics['avg_planned_time'] * 100
        
        # Filter routes with sufficient data
        route_metrics = route_metrics[route_metrics['total_count'] >= 5]  # Minimum 5 orders per route
        
        # Top routes by on-time rate
        top_on_time = route_metrics.sort_values('on_time_rate', ascending=False).head(top_n)
        
        # Bottom routes by on-time rate
        bottom_on_time = route_metrics.sort_values('on_time_rate').head(top_n)
        
        # Routes with highest average delay
        highest_delay = route_metrics.sort_values('avg_delay_time', ascending=False).head(top_n)
        
        result = {
            "top_on_time_routes": top_on_time[['route', 'on_time_rate', 'total_count', 'avg_actual_time', 'avg_planned_time']].to_dict('records'),
            "bottom_on_time_routes": bottom_on_time[['route', 'on_time_rate', 'total_count', 'avg_actual_time', 'avg_planned_time']].to_dict('records'),
            "highest_delay_routes": highest_delay[['route', 'avg_delay_time', 'on_time_rate', 'total_count']].to_dict('records'),
            "summary": {
                "total_routes_analyzed": len(route_metrics),
                "average_on_time_rate": round(float(route_metrics['on_time_rate'].mean()), 2),
                "average_time_efficiency": round(float(route_metrics['time_efficiency'].mean()), 2)
            }
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error analyzing key routes: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def analyze_province_route_times(file_path: str, origin_province_col: str, dest_province_col: str, actual_time_col: str) -> str:
    """
    Analyze average delivery times by origin and destination provinces

    Args:
        file_path (str): Path to CSV file
        origin_province_col (str): Column name for origin province
        dest_province_col (str): Column name for destination province
        actual_time_col (str): Column name for actual delivery time

    Returns:
        str: JSON string with province route time analysis
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        for col in [origin_province_col, dest_province_col, actual_time_col]:
            if col not in df.columns:
                return json.dumps({"error": f"Column not found: {col}"}, ensure_ascii=False)
        
        # Handle missing values
        df = df.dropna(subset=[origin_province_col, dest_province_col, actual_time_col])
        
        # Convert to numeric if needed
        df[actual_time_col] = pd.to_numeric(df[actual_time_col], errors='coerce')
        df = df.dropna(subset=[actual_time_col])
        
        # Calculate average times by origin-destination pair
        route_times = df.groupby([origin_province_col, dest_province_col]).agg({
            actual_time_col: ['mean', 'count', 'std']
        }).reset_index()
        
        # Flatten columns
        route_times.columns = ['origin_province', 'dest_province', 'avg_time', 'count', 'std_time']
        
        # Filter routes with sufficient data
        route_times = route_times[route_times['count'] >= 3]
        
        # Calculate average times by origin province
        origin_avg = df.groupby(origin_province_col)[actual_time_col].agg(['mean', 'count']).reset_index()
        origin_avg.columns = ['province', 'avg_time', 'count']
        origin_avg = origin_avg.sort_values('avg_time')
        
        # Calculate average times by destination province
        dest_avg = df.groupby(dest_province_col)[actual_time_col].agg(['mean', 'count']).reset_index()
        dest_avg.columns = ['province', 'avg_time', 'count']
        dest_avg = dest_avg.sort_values('avg_time')
        
        # Top 10 fastest routes
        fastest_routes = route_times.sort_values('avg_time').head(10)
        
        # Top 10 slowest routes
        slowest_routes = route_times.sort_values('avg_time', ascending=False).head(10)
        
        result = {
            "fastest_routes": fastest_routes[['origin_province', 'dest_province', 'avg_time', 'count']].to_dict('records'),
            "slowest_routes": slowest_routes[['origin_province', 'dest_province', 'avg_time', 'count']].to_dict('records'),
            "origin_province_averages": origin_avg.to_dict('records'),
            "dest_province_averages": dest_avg.to_dict('records'),
            "summary": {
                "total_routes_analyzed": len(route_times),
                "overall_average_time": round(float(df[actual_time_col].mean()), 2)
            }
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error analyzing province route times: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def analyze_process_times(file_path: str, process_columns: str) -> str:
    """
    Analyze times for different process stages

    Args:
        file_path (str): Path to CSV file
        process_columns (str): Comma-separated list of column names for different process stages

    Returns:
        str: JSON string with process time analysis
    """
    try:
        import pandas as pd
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Get process columns
        process_cols = [col.strip() for col in process_columns.split(',')]
        
        # Check if all columns exist
        missing_cols = [col for col in process_cols if col not in df.columns]
        if missing_cols:
            return json.dumps({"error": f"Columns not found: {', '.join(missing_cols)}"}, ensure_ascii=False)
        
        # Handle missing values and convert to numeric
        for col in process_cols:
            df[col] = pd.to_numeric(df[col], errors='coerce')
        
        # Remove rows with all missing process times
        df = df.dropna(subset=process_cols, how='all')
        
        # Calculate process time statistics
        process_stats = {}
        for col in process_cols:
            col_data = df[col].dropna()
            if len(col_data) > 0:
                process_stats[col] = {
                    "mean": round(float(col_data.mean()), 2),
                    "median": round(float(col_data.median()), 2),
                    "std": round(float(col_data.std()), 2),
                    "min": round(float(col_data.min()), 2),
                    "max": round(float(col_data.max()), 2),
                    "count": int(len(col_data)),
                    "missing": int(df[col].isnull().sum())
                }
            else:
                process_stats[col] = {
                    "mean": 0,
                    "median": 0,
                    "std": 0,
                    "min": 0,
                    "max": 0,
                    "count": 0,
                    "missing": int(len(df))
                }
        
        # Calculate total process time if multiple stages
        if len(process_cols) > 1:
            df['total_process_time'] = df[process_cols].sum(axis=1)
            total_data = df['total_process_time'].dropna()
            if len(total_data) > 0:
                process_stats['total_process_time'] = {
                    "mean": round(float(total_data.mean()), 2),
                    "median": round(float(total_data.median()), 2),
                    "std": round(float(total_data.std()), 2),
                    "min": round(float(total_data.min()), 2),
                    "max": round(float(total_data.max()), 2),
                    "count": int(len(total_data)),
                    "missing": int(df['total_process_time'].isnull().sum())
                }
        
        # Analyze process time distribution by stage
        stage_distribution = []
        if len(process_cols) > 1 and 'total_process_time' in df.columns:
            for col in process_cols:
                # Calculate percentage of total process time
                valid_rows = df.dropna(subset=[col, 'total_process_time'])
                if len(valid_rows) > 0:
                    avg_percentage = (valid_rows[col] / valid_rows['total_process_time']).mean() * 100
                    stage_distribution.append({
                        "stage": col,
                        "average_percentage": round(float(avg_percentage), 2),
                        "count": len(valid_rows)
                    })
        
        result = {
            "process_stage_statistics": process_stats,
            "stage_time_distribution": stage_distribution,
            "summary": {
                "total_records": len(df),
                "process_stages_analyzed": len(process_cols)
            }
        }
        
        return json.dumps(result, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error analyzing process times: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def generate_bar_chart(file_path: str, x_col: str, y_col: str, title: str = "Bar Chart", x_label: str = None, y_label: str = None, color: str = "#4CAF50", figsize: str = "10,6", save_path: str = None) -> str:
    """
    Generate bar chart from CSV data using matplotlib

    Args:
        file_path (str): Path to CSV file
        x_col (str): Column name for x-axis
        y_col (str): Column name for y-axis
        title (str): Chart title
        x_label (str): X-axis label (default: x_col)
        y_label (str): Y-axis label (default: y_col)
        color (str): Bar color (default: #4CAF50)
        figsize (str): Figure size as "width,height" (default: "10,6")
        save_path (str): Path to save the chart image (default: auto-generated)

    Returns:
        str: Path to the saved chart image
    """
    try:
        import pandas as pd
        import matplotlib.pyplot as plt
        import os
        from pathlib import Path
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        if x_col not in df.columns:
            return json.dumps({"error": f"Column not found: {x_col}"}, ensure_ascii=False)
        if y_col not in df.columns:
            return json.dumps({"error": f"Column not found: {y_col}"}, ensure_ascii=False)
        
        # Parse figsize
        try:
            width, height = map(float, figsize.split(','))
        except:
            width, height = 10, 6
        
        # Set default labels
        if x_label is None:
            x_label = x_col
        if y_label is None:
            y_label = y_col
        
        # Aggregate data
        if pd.api.types.is_numeric_dtype(df[y_col]):
            grouped = df.groupby(x_col)[y_col].mean().reset_index()
        else:
            # For categorical data
            grouped = df[x_col].value_counts().reset_index()
            grouped.columns = [x_col, y_col]
        
        # Create figure
        plt.figure(figsize=(width, height))
        
        # Generate bar chart
        bars = plt.bar(grouped[x_col].astype(str), grouped[y_col], color=color)
        
        # Add title and labels
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_label, fontsize=12)
        plt.ylabel(y_label, fontsize=12)
        
        # Rotate x-axis labels if needed
        if len(grouped) > 10:
            plt.xticks(rotation=45, ha='right')
        
        # Add value labels on bars
        for bar in bars:
            height_val = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height_val,
                    f'{height_val:.2f}',
                    ha='center', va='bottom')
        
        # Adjust layout
        plt.tight_layout()
        
        # Generate save path if not provided
        if save_path is None:
            # Create charts directory if it doesn't exist
            charts_dir = os.path.join(os.path.dirname(file_path), 'charts')
            os.makedirs(charts_dir, exist_ok=True)
            
            # Generate unique filename
            import uuid
            save_path = os.path.join(charts_dir, f'bar_chart_{uuid.uuid4().hex[:8]}.png')
        else:
            # Ensure directory exists
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
        
        # Save the chart
        plt.savefig(save_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return json.dumps({"chart_path": save_path, "title": title}, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error generating bar chart: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def generate_line_chart(file_path: str, date_col: str, value_col: str, title: str = "Line Chart", period: str = "daily", x_label: str = None, y_label: str = None, color: str = "#1976D2", figsize: str = "10,6", save_path: str = None) -> str:
    """
    Generate line chart from CSV data using matplotlib

    Args:
        file_path (str): Path to CSV file
        date_col (str): Column name for date/time
        value_col (str): Column name for value
        title (str): Chart title
        period (str): Time period (daily, weekly, monthly)
        x_label (str): X-axis label (default: period)
        y_label (str): Y-axis label (default: value_col)
        color (str): Line color (default: #1976D2)
        figsize (str): Figure size as "width,height" (default: "10,6")
        save_path (str): Path to save the chart image (default: auto-generated)

    Returns:
        str: Path to the saved chart image
    """
    try:
        import pandas as pd
        import matplotlib.pyplot as plt
        import os
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        if date_col not in df.columns:
            return json.dumps({"error": f"Column not found: {date_col}"}, ensure_ascii=False)
        if value_col not in df.columns:
            return json.dumps({"error": f"Column not found: {value_col}"}, ensure_ascii=False)
        
        # Convert date column to datetime
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        
        # Group by period
        if period == "daily":
            df['period'] = df[date_col].dt.date
        elif period == "weekly":
            df['period'] = df[date_col].dt.to_period('W')
        elif period == "monthly":
            df['period'] = df[date_col].dt.to_period('M')
        else:
            return json.dumps({"error": "Invalid period. Use daily, weekly, or monthly"}, ensure_ascii=False)
        
        # Parse figsize
        try:
            width, height = map(float, figsize.split(','))
        except:
            width, height = 10, 6
        
        # Set default labels
        if x_label is None:
            x_label = period
        if y_label is None:
            y_label = value_col
        
        # Aggregate data
        grouped = df.groupby('period')[value_col].mean().reset_index()
        grouped['period'] = grouped['period'].astype(str)
        
        # Create figure
        plt.figure(figsize=(width, height))
        
        # Generate line chart
        plt.plot(grouped['period'], grouped[value_col], marker='o', color=color, linewidth=2)
        
        # Add title and labels
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_label, fontsize=12)
        plt.ylabel(y_label, fontsize=12)
        
        # Rotate x-axis labels if needed
        if len(grouped) > 10:
            plt.xticks(rotation=45, ha='right')
        
        # Add grid
        plt.grid(True, linestyle='--', alpha=0.7)
        
        # Adjust layout
        plt.tight_layout()
        
        # Generate save path if not provided
        if save_path is None:
            # Create charts directory if it doesn't exist
            charts_dir = os.path.join(os.path.dirname(file_path), 'charts')
            os.makedirs(charts_dir, exist_ok=True)
            
            # Generate unique filename
            import uuid
            save_path = os.path.join(charts_dir, f'line_chart_{uuid.uuid4().hex[:8]}.png')
        else:
            # Ensure directory exists
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
        
        # Save the chart
        plt.savefig(save_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return json.dumps({"chart_path": save_path, "title": title}, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error generating line chart: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def add_web_link(url: str, title: str = None, categories: list = None) -> str:
    """
    Add a web link to the knowledge base

    Args:
        url (str): URL of the web page
        title (str): Title of the web page (optional)
        categories (list): List of categories for the link (optional)

    Returns:
        str: JSON string with link information
    """
    try:
        import requests
        from bs4 import BeautifulSoup
        from src.agents.knowledge_agent import KnowledgeStore
        
        # Fetch web page content
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        
        # Parse HTML to get title if not provided
        soup = BeautifulSoup(response.text, 'html.parser')
        if not title:
            title = soup.title.string if soup.title else url
        
        # Extract main content for summary
        content = ' '.join([p.get_text() for p in soup.find_all('p')[:10]])  # Get first 10 paragraphs
        
        # Generate summary using LLM
        from src.llms.llm import get_llm_by_name
        llm = get_llm_by_name("react_agent")
        summary_prompt = f"请为以下内容生成一个简洁的摘要（不超过200字）：\n{content}"
        summary = llm.invoke(summary_prompt).content.strip()
        
        # Add to knowledge store
        store = KnowledgeStore()
        link_data = store.add_link(url, title, summary, categories or [])
        
        return json.dumps({
            "success": True,
            "link": link_data
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error adding web link: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def get_link_summary(link_id: str) -> str:
    """
    Get the summary of a web link

    Args:
        link_id (str): ID of the link

    Returns:
        str: JSON string with link summary
    """
    try:
        from src.agents.knowledge_agent import KnowledgeStore
        store = KnowledgeStore()
        link = store.get_link(link_id)
        
        if not link:
            return json.dumps({"error": "Link not found"}, ensure_ascii=False)
        
        return json.dumps({
            "success": True,
            "link": link
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error getting link summary: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def search_links(query: str, categories: list = None) -> str:
    """
    Search links by query and optional categories

    Args:
        query (str): Search query
        categories (list): List of categories to filter by (optional)

    Returns:
        str: JSON string with search results
    """
    try:
        from src.agents.knowledge_agent import KnowledgeStore
        store = KnowledgeStore()
        results = store.search_links(query, categories)
        
        return json.dumps({
            "success": True,
            "results": results,
            "count": len(results)
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error searching links: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def categorize_link(link_id: str, categories: list) -> str:
    """
    Update the categories of a link

    Args:
        link_id (str): ID of the link
        categories (list): List of categories

    Returns:
        str: JSON string with updated link information
    """
    try:
        from src.agents.knowledge_agent import KnowledgeStore
        store = KnowledgeStore()
        link = store.categorize_link(link_id, categories)
        
        if not link:
            return json.dumps({"error": "Link not found"}, ensure_ascii=False)
        
        return json.dumps({
            "success": True,
            "link": link
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error categorizing link: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def list_links(categories: list = None) -> str:
    """
    List all links, optionally filtered by categories

    Args:
        categories (list): List of categories to filter by (optional)

    Returns:
        str: JSON string with list of links
    """
    try:
        from src.agents.knowledge_agent import KnowledgeStore
        store = KnowledgeStore()
        links = store.list_links(categories)
        
        return json.dumps({
            "success": True,
            "links": links,
            "count": len(links)
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error listing links: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def delete_link(link_id: str) -> str:
    """
    Delete a link by its ID

    Args:
        link_id (str): ID of the link

    Returns:
        str: JSON string with deletion status
    """
    try:
        from src.agents.knowledge_agent import KnowledgeStore
        store = KnowledgeStore()
        success = store.delete_link(link_id)
        
        if not success:
            return json.dumps({"error": "Link not found"}, ensure_ascii=False)
        
        return json.dumps({
            "success": True,
            "message": "Link deleted successfully"
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error deleting link: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def list_documents(directory: str = None) -> str:
    """
    List all documents in the specified directory

    Args:
        directory (str): Directory path to list documents from (default: documents directory)

    Returns:
        str: JSON string containing array of document information
    """
    try:
        # Default documents directory
        if directory is None:
            directory = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "documents")
        
        # Create directory if it doesn't exist
        os.makedirs(directory, exist_ok=True)
        
        # Scan directory for documents
        path = Path(directory)
        supported_extensions = ['.docx', '.txt', '.xlsx', '.xls', '.ppt', '.pptx', '.md', '.py', '.csv']
        documents = []
        
        for ext in supported_extensions:
            documents.extend(path.glob(f'*{ext}'))
        
        docs_info = []
        for doc_path in documents:
            try:
                file_info = {
                    "file_name": doc_path.name,
                    "file_path": str(doc_path).replace('\\', '/'),
                    "file_type": doc_path.suffix.lower(),
                    "size": doc_path.stat().st_size,
                    "modified": doc_path.stat().st_mtime
                }
                docs_info.append(file_info)
            except Exception as e:
                logger.warning(f"Failed to process document {doc_path}: {e}")
                continue
        
        return json.dumps({
            "directory": str(directory).replace('\\', '/'),
            "documents": docs_info
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Failed to list documents: {e}", exc_info=True)
        return json.dumps({
            "error": f"Failed to list documents: {str(e)}",
            "documents": []
        }, ensure_ascii=False, indent=2)


@tool
def upload_document(file_content: str, file_name: str, directory: str = None) -> str:
    """
    Upload a document to the specified directory

    Args:
        file_content (str): Base64 encoded file content
        file_name (str): Name of the file to upload
        directory (str): Directory path to upload to (default: documents directory)

    Returns:
        str: JSON string with upload status
    """
    try:
        import base64
        
        # Default documents directory
        if directory is None:
            directory = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "documents")
        
        # Create directory if it doesn't exist
        os.makedirs(directory, exist_ok=True)
        
        # Decode base64 content
        try:
            decoded_content = base64.b64decode(file_content)
        except Exception as e:
            return json.dumps({"error": f"Failed to decode file content: {str(e)}"}, ensure_ascii=False)
        
        # Save file
        file_path = os.path.join(directory, file_name)
        with open(file_path, "wb") as f:
            f.write(decoded_content)
        
        return json.dumps({
            "success": True,
            "file_path": file_path.replace('\\', '/'),
            "message": f"Document {file_name} uploaded successfully"
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Failed to upload document: {e}", exc_info=True)
        return json.dumps({
            "error": f"Failed to upload document: {str(e)}"
        }, ensure_ascii=False, indent=2)


@tool
def delete_document(file_path: str) -> str:
    """
    Delete a document

    Args:
        file_path (str): Path to the document to delete

    Returns:
        str: JSON string with deletion status
    """
    try:
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Delete the file
        os.remove(file_path)
        
        return json.dumps({
            "success": True,
            "message": f"Document {os.path.basename(file_path)} deleted successfully"
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Failed to delete document: {e}", exc_info=True)
        return json.dumps({
            "error": f"Failed to delete document: {str(e)}"
        }, ensure_ascii=False, indent=2)


@tool
def preview_document(file_path: str, max_lines: int = 100) -> str:
    """
    Preview a document's content

    Args:
        file_path (str): Path to the document to preview
        max_lines (int): Maximum number of lines to return (default: 100)

    Returns:
        str: JSON string with document preview
    """
    try:
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        file_ext = Path(file_path).suffix.lower()
        content = ""
        
        if file_ext == '.txt' or file_ext == '.md' or file_ext == '.py':
            # Read text files directly
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = f.readlines()
                content = ''.join(lines[:max_lines])
                if len(lines) > max_lines:
                    content += f"\n... (truncated, total {len(lines)} lines)"
        
        elif file_ext == '.csv':
            # Read CSV files with pandas
            import pandas as pd
            df = pd.read_csv(file_path)
            content = df.head(min(20, len(df))).to_string()
            if len(df) > 20:
                content += f"\n... (truncated, total {len(df)} rows)"
        
        elif file_ext in ['.xlsx', '.xls']:
            # Read Excel files with pandas
            import pandas as pd
            df = pd.read_excel(file_path)
            content = df.head(min(20, len(df))).to_string()
            if len(df) > 20:
                content += f"\n... (truncated, total {len(df)} rows)"
        
        elif file_ext == '.docx':
            # Read Word documents
            try:
                import docx
                doc = docx.Document(file_path)
                paragraphs = []
                for para in doc.paragraphs[:max_lines]:
                    paragraphs.append(para.text)
                content = '\n'.join(paragraphs)
                if len(doc.paragraphs) > max_lines:
                    content += f"\n... (truncated, total {len(doc.paragraphs)} paragraphs)"
            except ImportError:
                content = "Preview not available: python-docx library not installed"
            except Exception as e:
                content = f"Preview error: {str(e)}"
        
        elif file_ext in ['.ppt', '.pptx']:
            # Read PowerPoint presentations
            try:
                import pptx
                presentation = pptx.Presentation(file_path)
                slides = []
                for i, slide in enumerate(presentation.slides[:10]):  # Limit to first 10 slides
                    slide_text = f"Slide {i+1}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            slide_text += f"  {shape.text}\n"
                    slides.append(slide_text)
                content = '\n'.join(slides)
                if len(presentation.slides) > 10:
                    content += f"\n... (truncated, total {len(presentation.slides)} slides)"
            except ImportError:
                content = "Preview not available: python-pptx library not installed"
            except Exception as e:
                content = f"Preview error: {str(e)}"
        
        else:
            content = f"Preview not supported for {file_ext} files"
        
        return json.dumps({
            "file_path": file_path,
            "file_type": file_ext,
            "content": content
        }, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Failed to preview document: {e}", exc_info=True)
        return json.dumps({
            "error": f"Failed to preview document: {str(e)}"
        }, ensure_ascii=False, indent=2)


@tool
def generate_pie_chart(file_path: str, category_col: str, title: str = "Pie Chart", figsize: str = "8,8", colors: str = None, explode: str = None, save_path: str = None) -> str:
    """
    Generate pie chart from CSV data using matplotlib

    Args:
        file_path (str): Path to CSV file
        category_col (str): Column name for categories
        title (str): Chart title
        figsize (str): Figure size as "width,height" (default: "8,8")
        colors (str): Comma-separated list of colors (default: matplotlib default)
        explode (str): Comma-separated list of explode values (default: None)
        save_path (str): Path to save the chart image (default: auto-generated)

    Returns:
        str: Path to the saved chart image
    """
    try:
        import pandas as pd
        import matplotlib.pyplot as plt
        import os
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required column exists
        if category_col not in df.columns:
            return json.dumps({"error": f"Column not found: {category_col}"}, ensure_ascii=False)
        
        # Parse figsize
        try:
            width, height = map(float, figsize.split(','))
        except:
            width, height = 8, 8
        
        # Count occurrences
        counts = df[category_col].value_counts()
        labels = counts.index.astype(str)
        values = counts.values
        
        # Parse colors
        color_list = None
        if colors:
            color_list = [c.strip() for c in colors.split(',')]
        
        # Parse explode
        explode_list = None
        if explode:
            try:
                explode_list = [float(e.strip()) for e in explode.split(',')]
                # Ensure explode list matches length of labels
                if len(explode_list) != len(labels):
                    explode_list = None
            except:
                explode_list = None
        
        # Create figure
        plt.figure(figsize=(width, height))
        
        # Generate pie chart
        wedges, texts, autotexts = plt.pie(
            values,
            labels=labels,
            autopct='%1.1f%%',
            startangle=90,
            colors=color_list,
            explode=explode_list,
            shadow=True
        )
        
        # Customize text
        for text in texts:
            text.set_fontsize(12)
        for autotext in autotexts:
            autotext.set_fontsize(10)
            autotext.set_color('white')
            autotext.set_fontweight('bold')
        
        # Add title
        plt.title(title, fontsize=16, fontweight='bold')
        
        # Add legend
        plt.legend(wedges, labels, title="Categories", loc="best")
        
        # Ensure equal aspect ratio
        plt.axis('equal')
        
        # Adjust layout
        plt.tight_layout()
        
        # Generate save path if not provided
        if save_path is None:
            # Create charts directory if it doesn't exist
            charts_dir = os.path.join(os.path.dirname(file_path), 'charts')
            os.makedirs(charts_dir, exist_ok=True)
            
            # Generate unique filename
            import uuid
            save_path = os.path.join(charts_dir, f'pie_chart_{uuid.uuid4().hex[:8]}.png')
        else:
            # Ensure directory exists
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
        
        # Save the chart
        plt.savefig(save_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return json.dumps({"chart_path": save_path, "title": title}, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error generating pie chart: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


@tool
def generate_scatter_plot(file_path: str, x_col: str, y_col: str, title: str = "Scatter Plot", x_label: str = None, y_label: str = None, color: str = "#FF5722", figsize: str = "10,6", alpha: float = 0.6, save_path: str = None) -> str:
    """
    Generate scatter plot from CSV data using matplotlib

    Args:
        file_path (str): Path to CSV file
        x_col (str): Column name for x-axis
        y_col (str): Column name for y-axis
        title (str): Chart title
        x_label (str): X-axis label (default: x_col)
        y_label (str): Y-axis label (default: y_col)
        color (str): Point color (default: #FF5722)
        figsize (str): Figure size as "width,height" (default: "10,6")
        alpha (float): Point transparency (default: 0.6)
        save_path (str): Path to save the chart image (default: auto-generated)

    Returns:
        str: Path to the saved chart image
    """
    try:
        import pandas as pd
        import matplotlib.pyplot as plt
        import os
        
        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}"}, ensure_ascii=False)
        
        # Try different separators
        separators = ['^', ',', '\t']
        df = None
        
        for sep in separators:
            try:
                df = pd.read_csv(file_path, sep=sep)
                if len(df.columns) > 1:
                    break
            except:
                continue
        
        if df is None:
            return json.dumps({"error": "Failed to read CSV file"}, ensure_ascii=False)
        
        # Check if required columns exist
        if x_col not in df.columns:
            return json.dumps({"error": f"Column not found: {x_col}"}, ensure_ascii=False)
        if y_col not in df.columns:
            return json.dumps({"error": f"Column not found: {y_col}"}, ensure_ascii=False)
        
        # Parse figsize
        try:
            width, height = map(float, figsize.split(','))
        except:
            width, height = 10, 6
        
        # Set default labels
        if x_label is None:
            x_label = x_col
        if y_label is None:
            y_label = y_col
        
        # Filter numeric data
        numeric_df = df[[x_col, y_col]].apply(pd.to_numeric, errors='coerce').dropna()
        
        # Limit to 1000 points
        if len(numeric_df) > 1000:
            numeric_df = numeric_df.sample(1000)
        
        # Create figure
        plt.figure(figsize=(width, height))
        
        # Generate scatter plot
        plt.scatter(numeric_df[x_col], numeric_df[y_col], color=color, alpha=alpha, edgecolors='black', s=50)
        
        # Add title and labels
        plt.title(title, fontsize=16, fontweight='bold')
        plt.xlabel(x_label, fontsize=12)
        plt.ylabel(y_label, fontsize=12)
        
        # Add grid
        plt.grid(True, linestyle='--', alpha=0.7)
        
        # Adjust layout
        plt.tight_layout()
        
        # Generate save path if not provided
        if save_path is None:
            # Create charts directory if it doesn't exist
            charts_dir = os.path.join(os.path.dirname(file_path), 'charts')
            os.makedirs(charts_dir, exist_ok=True)
            
            # Generate unique filename
            import uuid
            save_path = os.path.join(charts_dir, f'scatter_plot_{uuid.uuid4().hex[:8]}.png')
        else:
            # Ensure directory exists
            os.makedirs(os.path.dirname(save_path), exist_ok=True)
        
        # Save the chart
        plt.savefig(save_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return json.dumps({"chart_path": save_path, "title": title}, ensure_ascii=False, indent=2)
        
    except Exception as e:
        logger.error(f"Error generating scatter plot: {e}")
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def read_file(
        file_path: str,
        mode: Literal["all", "head", "tail"] = "all",
        n_lines: Optional[int] = None
) -> str:
    """
    Read file content
    * Avoid reading large files that produce excessive output.

    Args:
        file_path (str): File path
        mode (str): Read mode, options:
            - "all": Read entire file (default)
            - "head": Read first N lines
            - "tail": Read last N lines
        n_lines (int, optional): Number of lines to read when mode is "head" or "tail"

    Returns:
        str: File content or error message
    """
    try:
        # Check if file exists
        if not os.path.exists(file_path):
            return f"Error: File '{file_path}' does not exist"

        # Check if it is a file
        if not os.path.isfile(file_path):
            return f"Error: '{file_path}' is not a file"

        # Read file
        with open(file_path, 'r', encoding='utf-8') as f:
            if mode == "all":
                # Read entire file
                try:
                    # Prioritize file size limit to avoid loading large files entirely into memory
                    file_size = os.path.getsize(file_path)
                    if file_size > _MAX_RETURN_BYTES:
                        raise ValueError(f"Read result is too large, exceeding token limit")
                except Exception:
                    # If getting file size fails, ignore the error here and check content size later
                    pass

                content = f.read()
                _ensure_within_limit_by_bytes(content, context=f"all:{file_path}")
                return content

            elif mode == "head":
                # Read first N lines
                if n_lines is None:
                    return "Error: n_lines parameter must be specified when using head mode"

                lines = []
                for i, line in enumerate(f):
                    if i >= n_lines:
                        break
                    lines.append(line)
                content = ''.join(lines)
                _ensure_within_limit_by_bytes(content, context=f"head(n={n_lines}):{file_path}")
                return content

            elif mode == "tail":
                # Read last N lines
                if n_lines is None:
                    return "Error: n_lines parameter must be specified when using tail mode"

                # Read all lines and keep last N lines
                all_lines = f.readlines()
                tail_lines = all_lines[-n_lines:] if len(all_lines) >= n_lines else all_lines
                content = ''.join(tail_lines)
                _ensure_within_limit_by_bytes(content, context=f"tail(n={n_lines}):{file_path}")
                return content

            else:
                return f"Error: Unsupported mode '{mode}', please use 'all', 'head' or 'tail'"

    except UnicodeDecodeError:
        # If UTF-8 decoding fails, try other encodings
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                if mode == "all":
                    content = f.read()
                    _ensure_within_limit_by_bytes(content, context=f"all(gbk):{file_path}")
                    return content
                elif mode == "head":
                    content = ''.join([f.readline() for _ in range(n_lines or 0)])
                    _ensure_within_limit_by_bytes(content, context=f"head(gbk,n={n_lines}):{file_path}")
                    return content
                elif mode == "tail":
                    all_lines = f.readlines()
                    content = ''.join(all_lines[-n_lines:] if len(all_lines) >= n_lines else all_lines)
                    _ensure_within_limit_by_bytes(content, context=f"tail(gbk,n={n_lines}):{file_path}")
                    return content
                return f"Error: Unsupported mode '{mode}', please use 'all', 'head' or 'tail'"
        except Exception as e:
            return f"Error: Unable to read file, encoding issue - {str(e)}"

    except Exception as e:
        return f"Error: Exception occurred while reading file - {str(e)}"

